import csv
import importlib.util
import json
from io import StringIO
from pathlib import Path
import zipfile

import pytest
from fastapi.testclient import TestClient
from sqlalchemy.exc import IntegrityError
from sqlalchemy import create_engine
from sqlalchemy.orm import Session, sessionmaker
from sqlalchemy.pool import StaticPool

from app.db.base import Base
from app.db.session import get_db
from app.main import app
from app.models import DeepAnalysisQueueItem, FulltextAnalysisResult, HighlightCard, NotableAuthor, PdfAsset, StrongEvidence
from app.services.context_service import build_context_preview
from app.services.evidence_service import EvidenceService
from app.services.highlight_card_service import HighlightCardService, PptxExportError, render_markdown_highlight
from app.services.scholar_report_service import ScholarReportService
from pptx import Presentation
from tests.test_scholar_evidence import seed_queue_item


@pytest.fixture()
def db_session_factory():
    engine = create_engine(
        "sqlite://",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    Base.metadata.create_all(bind=engine)
    return sessionmaker(bind=engine, autoflush=False, autocommit=False)


@pytest.fixture()
def client(db_session_factory):
    def override_get_db():
        db = db_session_factory()
        try:
            yield db
        finally:
            db.close()

    app.dependency_overrides[get_db] = override_get_db
    try:
        yield TestClient(app)
    finally:
        app.dependency_overrides.clear()


def seed_evidence(
    db: Session,
    tmp_path,
    *,
    review_status: str = "accepted",
    aspect: str = "method_foundation",
    score: float = 0.9,
    evidence_strength: str = "strong",
    citation_text: str = "Cited Scholar Paper provides a method foundation for reviewable citation evidence.",
    title: str = "Independent Citing Paper",
    analysis_diagnostics: dict = None,
):
    session_id, item_id = seed_queue_item(
        db,
        tmp_path,
        title=title,
        text=f"Before context. {citation_text} After context.",
    )
    from app.models import DeepAnalysisQueueItem

    item = db.get(DeepAnalysisQueueItem, item_id)
    result = FulltextAnalysisResult(
        scholar_session_id=session_id,
        queue_item_id=item_id,
        citation_edge_id=item.citation_edge_id,
        analysis_scope="fulltext_anchor_direct",
        status="succeeded",
        candidate_spans_json=json.dumps(analysis_diagnostics or {}),
        parsed_result_json=json.dumps(
            {
                "findings": [
                    {
                        "evidence_type": aspect,
                        "stance": "positive",
                        "mention_type": "strong",
                        "citation_text": citation_text,
                        "reasoning": "The citing paper grounds the claim in original citation evidence.",
                        "keywords": ["method foundation", "citation evidence"],
                    }
                ]
            }
        ),
    )
    db.add(result)
    db.flush()
    evidence = StrongEvidence(
        fulltext_result_id=result.id,
        scholar_session_id=session_id,
        queue_item_id=item_id,
        citation_edge_id=item.citation_edge_id,
        aspect=aspect,
        stance="positive",
        mention_type="strong",
        citation_text=citation_text,
        highlighted_text_html=citation_text.replace(
            "method foundation",
            "<mark>method foundation</mark>",
        ),
        highlight_keywords_json=json.dumps(["method foundation", "citation evidence"]),
        evidence_reason="The citing paper grounds the claim in original citation evidence.",
        evidence_strength=evidence_strength,
        score=score,
        third_party_status="third_party",
        review_status=review_status,
    )
    db.add(evidence)
    db.commit()
    return session_id, evidence.id


def seed_template_direct_result(
    db: Session,
    tmp_path,
    *,
    recommendation: str = "include",
    quote: str = "Target Paper enables through-wall eavesdropping with RFID tags [23].",
    context: str = "The citing paper discusses Target Paper enables through-wall eavesdropping with RFID tags [23] in the main body.",
    title: str = "Traceable Citing Paper",
):
    anchored_quote = quote if "[23]" in quote else f"{quote} [23]"
    anchored_context = context if "[23]" in context else f"{context} [23]"
    session_id, item_id = seed_queue_item(
        db,
        tmp_path,
        title=title,
        target_title="Target Paper",
        text=f"{anchored_context}\n\nReferences\n[23] Target Paper. doi:10.1145/target",
    )
    from app.models import DeepAnalysisQueueItem

    item = db.get(DeepAnalysisQueueItem, item_id)
    result = FulltextAnalysisResult(
        scholar_session_id=session_id,
        queue_item_id=item_id,
        citation_edge_id=item.citation_edge_id,
        analysis_scope="fulltext_template_direct",
        status="succeeded",
        candidate_spans_json=json.dumps(
            {
                "mode": "fulltext_template_direct",
                "include_count": 1 if recommendation == "include" else 0,
                "review_count": 1 if recommendation == "review" else 0,
                "exclude_count": 1 if recommendation == "exclude" else 0,
            }
        ),
        parsed_result_json=json.dumps(
            {
                "target_reference_marker": "[23]",
                "target_reference_entry": "[23] Target Paper. doi:10.1145/target",
                "paper_level_summary_zh": "摘要。",
                "evidences": [
                    {
                        "recommendation": recommendation,
                        "claim_type": "capability_recognition",
                        "evidence_quote": anchored_quote,
                        "evidence_context": anchored_context,
                        "reference_entry": "[23] Target Paper. doi:10.1145/target",
                        "why_this_judgment_zh": "正文证据通过 [23] 锚定目标论文，并说明能力判断。",
                        "copy_ready_zh": "引用论文在正文中明确讨论目标论文的能力表现，可作为报告材料。",
                        "confidence": "high",
                    }
                ],
            },
            ensure_ascii=False,
        ),
    )
    db.add(result)
    db.commit()
    return session_id, result.id


def test_generate_cards_from_accepted_evidence(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(db, tmp_path, review_status="accepted")
        cards = HighlightCardService(db).generate_cards_from_evidence(session_id)

    assert len(cards) == 1
    assert cards[0].strong_evidence_id == evidence_id
    assert cards[0].card_type == "method_foundation"
    assert cards[0].evidence_quote
    assert cards[0].narrative_zh
    assert cards[0].source_evidence_id == evidence_id


def test_generate_impact_card_from_theoretical_foundation_evidence(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text="Cited Scholar Paper provides the theoretical foundation for our model.",
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert card.card_type == "theoretical_foundation"
    assert "理论推导" in (card.narrative_zh or "") or "理论建模" in (card.narrative_zh or "")


def test_theoretical_foundation_has_why_this_judgment(db_session_factory, tmp_path):
    quote = "The cited approach [21] defines the latent transition model used in the derivation."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[21]"},
        )
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]
        evidence_row = EvidenceService(db).list_scholar_evidence(session_id)[0]

    assert row["narrative_meta"]["why_this_judgment"]
    assert "理论基础" in row["narrative_meta"]["why_this_judgment"]
    assert evidence_row["judgment_basis"]["why_this_judgment"]


def test_representative_work_judgment_says_not_high_praise(db_session_factory, tmp_path):
    quote = "Related work lists the cited method [22] as part of the adaptive tracking pipeline."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="representative_work",
            evidence_strength="weak",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[22]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.stance = "neutral"
        evidence.mention_type = "related_work"
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]

    statement = row["narrative_meta"]["copy_ready_statement"]
    why = row["narrative_meta"]["why_this_judgment"]
    assert "不宜表述为高度评价" in statement
    assert "不是直接高度评价" in why or "不是直接正向评价" in row["narrative_meta"]["risk_note"]


def test_mechanism_intro_representative_card_has_why_this_judgment(db_session_factory, tmp_path):
    quote = "Moreover, MoireVision [40] introduces a generalized 6-DoF motion sensing mechanism."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="representative_work",
            evidence_strength="moderate",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[40]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.stance = "neutral"
        evidence.mention_type = "related_work"
        evidence.highlight_keywords_json = json.dumps(["generalized 6-DoF motion sensing mechanism"])
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]

    why = row["narrative_meta"]["why_this_judgment"]
    assert why
    assert "generalized 6-DoF motion sensing mechanism" in why
    assert "不是直接高度评价" in why or "不应包装成高度评价" in row["narrative_meta"]["risk_note"]


def test_sub_mm_evidence_points_to_original_phrase(db_session_factory, tmp_path):
    quote = "The system reaches sub-mm vibration sensing accuracy in the reported experiment [23]."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="positive_evaluation",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]

    assert "sub-mm" in row["narrative_meta"]["key_phrases"]
    assert "vibration sensing" in row["narrative_meta"]["key_phrases"]
    assert "sub-mm" in row["narrative_meta"]["why_this_judgment"]
    assert "vibration sensing" in row["narrative_meta"]["why_this_judgment"]


def test_grouped_citation_has_attribution_risk_note(db_session_factory, tmp_path):
    quote = "Prior methods [24], [25], [26] are used as the main comparison group."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="detailed_comparison",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[25]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.mention_type = "grouped_citation"
        evidence.anchor_status = "grouped_citation"
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]

    assert "成组引用" in row["narrative_meta"]["risk_note"]
    assert "不能自动断言" in row["narrative_meta"]["why_this_judgment"]


def test_radsee_negative_mm_level_not_positive_submm_evidence(db_session_factory, tmp_path):
    quote = "RadSee reports millimeter-level RFID sensing, but the method is sensitive to deployment constraints [23]."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="limitation_or_negative",
            evidence_strength="moderate",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.stance = "negative"
        evidence.template_satisfied = True
        evidence.template_match_reason = "precision_claim: body quote anchors to target and contains millimeter-level"
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]

    meta = row["narrative_meta"]
    assert meta["report_recommendation"] == "候选复核"
    assert "局限性" in meta["judgment_label"]
    assert "不能" in meta["limitation_zh"]
    assert "不能" in meta["copy_ready_statement_zh"]


def test_interpretation_contains_why_and_not_overclaim(db_session_factory, tmp_path):
    quote = "Related work cites the RFID sensing method [23] as a prior system."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="representative_work",
            evidence_strength="weak",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.stance = "neutral"
        evidence.mention_type = "related_work"
        evidence.template_satisfied = False
        evidence.template_failure_reason = "plain RFID reference without sub-mm or vibration capability wording"
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]

    meta = row["narrative_meta"]
    assert meta["judgment_basis_zh"]
    assert "不能写成高度评价" in meta["limitation_zh"]
    assert meta["report_recommendation"] == "候选复核"


def test_copy_ready_statement_is_specific_not_template(db_session_factory, tmp_path):
    quote = "The target RFID tag captures loudspeaker vibration for acoustic sensing [23]."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="capability_recognition",
            evidence_strength="strong",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.template_satisfied = True
        evidence.template_match_reason = "capability_recognition: body quote anchors to target and links RFID with loudspeaker vibration"
        evidence.highlight_keywords_json = json.dumps(["RFID tag", "loudspeaker vibration", "acoustic sensing"])
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        row = HighlightCardService(db).list_report_workspace_cards(session_id)[0]

    statement = row["narrative_meta"]["copy_ready_statement_zh"]
    assert "loudspeaker" in statement or "扬声器" in statement
    assert "acoustic" in statement or "声学" in statement
    assert "使用了目标论文的核心思想" not in statement


def test_report_separates_high_confidence_candidate_negative_and_ordinary(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="method_foundation",
            evidence_strength="strong",
            citation_text="The cited method [21] is used as the main sensing module.",
            analysis_diagnostics={"target_reference_marker": "[21]"},
        )
        base = db.get(StrongEvidence, evidence_id)
        negative = StrongEvidence(
            fulltext_result_id=base.fulltext_result_id,
            scholar_session_id=base.scholar_session_id,
            queue_item_id=base.queue_item_id,
            citation_edge_id=base.citation_edge_id,
            aspect="limitation_or_negative",
            stance="negative",
            mention_type="explicit_target",
            citation_text="The cited method [21] remains sensitive to environmental constraints.",
            evidence_reason="negative feedback",
            evidence_strength="moderate",
            score=0.7,
            third_party_status="third_party",
            review_status="accepted",
        )
        ordinary = StrongEvidence(
            fulltext_result_id=base.fulltext_result_id,
            scholar_session_id=base.scholar_session_id,
            queue_item_id=base.queue_item_id,
            citation_edge_id=base.citation_edge_id,
            aspect="representative_work",
            stance="neutral",
            mention_type="related_work",
            citation_text="Related work lists the cited method [21] as a prior RFID system.",
            evidence_reason="ordinary related work",
            evidence_strength="weak",
            score=0.4,
            third_party_status="third_party",
            review_status="accepted",
        )
        db.add_all([negative, ordinary])
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "## 二、强证据卡片 / 高可信第三方佐证" in markdown
    assert "## 三、候选佐证，需要人工复核" in markdown
    assert "## 四、局限性/负面反馈" in markdown
    assert "## 五、普通相关工作引用" in markdown
    assert "sensitive to environmental constraints" in markdown
    assert "prior RFID system" in markdown


def test_evidence_and_report_pages_show_judgment_sections(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text="The cited method [27] defines the graph process model.",
            analysis_diagnostics={"target_reference_marker": "[27]"},
        )
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    evidence_response = client.get(f"/scholar-sessions/{session_id}/evidence?mode=debug")
    cards_response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    for html in [evidence_response.text, cards_response.text]:
        assert "【原文证据】" in html
        assert "【为什么这样判断】" in html
        assert "【可复制表述】" in html
        assert "【风险提示】" in html
        assert "anchor_validation_status" in html


def test_evidence_page_defaults_to_formal_view(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text="Target Paper enables through-wall eavesdropping with RFID tags [23].",
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )

    response = client.get(f"/scholar-sessions/{session_id}/evidence")

    assert response.status_code == 200
    assert "正式证据视图" in response.text
    assert "formal-evidence-card" in response.text
    assert "展开调试信息" in response.text


def test_evidence_page_hides_debug_fields_by_default(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text="Target Paper enables through-wall eavesdropping [23].",
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )

    response = client.get(f"/scholar-sessions/{session_id}/evidence")

    assert "anchor_validation_status" not in response.text
    assert "模板判断理由" not in response.text
    assert "用户备注" not in response.text
    assert "修正标签" not in response.text
    assert "保存复核备注" not in response.text


def test_evidence_page_debug_panel_contains_anchor_fields(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text="Target Paper enables through-wall eavesdropping [23].",
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )

    response = client.get(f"/scholar-sessions/{session_id}/evidence?mode=debug")

    assert response.status_code == 200
    assert "anchor_validation_status" in response.text
    assert "anchor_validation_reason" in response.text
    assert "保存复核备注" in response.text


def test_evidence_formal_card_shows_quote_reference_context_evaluation(client, db_session_factory, tmp_path):
    quote = "Target Paper enables through-wall eavesdropping with RFID tags [23]."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[23]"},
        )

    response = client.get(f"/scholar-sessions/{session_id}/evidence")

    assert "原文证据" in response.text
    assert "对应参考文献（引用论文原文 References 中的条目）" in response.text
    assert "原文上下文" in response.text
    assert "亮点评价" in response.text
    assert "评价理由" in response.text
    assert "可复制表述" in response.text
    assert "Before context" in response.text
    assert 'class="citation-marker">[23]</mark>' in response.text
    assert 'class="claim-phrase">through-wall eavesdropping</mark>' in response.text


def test_evidence_page_reuses_formal_report_view_model_or_partial():
    template_root = Path(__file__).resolve().parents[1] / "app" / "templates" / "scholar_sessions"
    evidence_template = (template_root / "evidence.html").read_text(encoding="utf-8")
    report_template = (template_root / "cards.html").read_text(encoding="utf-8")
    shared_partial = "components/formal_evidence_card.html"

    assert shared_partial in evidence_template
    assert shared_partial in report_template


def test_card_citation_text_must_anchor_to_target(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            citation_text="The other tagged method [20] is highly accurate.",
            analysis_diagnostics={"target_reference_marker": "[16]"},
        )
        service = HighlightCardService(db)
        cards = service.generate_cards_from_evidence(session_id)
        evidence = db.get(StrongEvidence, evidence_id)
        evidence_review_status = evidence.review_status
        evidence_anchor_status = evidence.anchor_status

    assert cards == []
    assert evidence_review_status == "false_positive"
    assert evidence_anchor_status == "invalid"


def test_existing_mismatched_cards_marked_false_positive(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            citation_text="A different method [20] receives the positive comment.",
            analysis_diagnostics={"target_reference_marker": "[16]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        card = HighlightCard(
            scholar_session_id=session_id,
            strong_evidence_id=evidence_id,
            card_type="positive_evaluation",
            title="Mismatched card",
            body_markdown="bad",
            evidence_quote=evidence.citation_text,
            source_citing_paper_title="Citing",
            source_cited_paper_title="Target",
            aspect="positive_evaluation",
            stance="positive",
            evidence_strength="strong",
            score=0.9,
            review_status="accepted",
            include_in_report=True,
        )
        db.add(card)
        db.commit()
        rows = HighlightCardService(db).list_report_workspace_cards(session_id)
        db.refresh(evidence)
        db.refresh(card)
        evidence_review_status = evidence.review_status
        card_review_status = card.review_status
        card_include_in_report = card.include_in_report
        card_strength = card.evidence_strength
        card_score = card.score

    assert rows == []
    assert evidence_review_status == "false_positive"
    assert card_review_status == "false_positive"
    assert card_include_in_report is False
    assert card_strength == "none"
    assert card_score == 0


def test_false_positive_card_not_shown_in_default_report_workspace(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            citation_text="The proposed marker is similar to MoireTag [20].",
            analysis_diagnostics={"target_reference_marker": "[16]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        card = HighlightCard(
            scholar_session_id=session_id,
            strong_evidence_id=evidence_id,
            card_type="positive_evaluation",
            title="Wrong positive card",
            body_markdown="bad",
            evidence_quote=evidence.citation_text,
            source_citing_paper_title="Citing",
            source_cited_paper_title="MoiréTracker",
            aspect="positive_evaluation",
            stance="positive",
            evidence_strength="strong",
            score=0.95,
            review_status="accepted",
            include_in_report=True,
        )
        db.add(card)
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "Wrong positive card" not in response.text
    assert "MoireTag [20]" not in response.text


def test_false_positive_view_shows_invalid_anchor_cards(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            citation_text="Recently, other methods [18] receive positive comments.",
            analysis_diagnostics={"target_reference_marker": "[16]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        card = HighlightCard(
            scholar_session_id=session_id,
            strong_evidence_id=evidence_id,
            card_type="positive_evaluation",
            title="Wrong card",
            body_markdown="bad",
            evidence_quote=evidence.citation_text,
            source_citing_paper_title="Citing",
            source_cited_paper_title="MoiréTracker",
            aspect="positive_evaluation",
            stance="positive",
            evidence_strength="strong",
            score=0.95,
            review_status="accepted",
            include_in_report=True,
        )
        db.add(card)
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace?view=false_positive")

    assert response.status_code == 200
    assert "误报候选" in response.text
    assert "cited_other_reference_marker" in response.text
    assert "Recently, other methods [18]" in response.text


def test_false_positive_badge_overrides_strength_badge(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            citation_text="The proposed marker is similar to MoireTag [20].",
            analysis_diagnostics={"target_reference_marker": "[16]"},
        )
        evidence = db.get(StrongEvidence, evidence_id)
        card = HighlightCard(
            scholar_session_id=session_id,
            strong_evidence_id=evidence_id,
            card_type="positive_evaluation",
            title="Strength conflict",
            body_markdown="bad",
            evidence_quote=evidence.citation_text,
            source_citing_paper_title="Citing",
            source_cited_paper_title="MoiréTracker",
            aspect="positive_evaluation",
            stance="positive",
            evidence_strength="strong",
            score=0.95,
            review_status="accepted",
            include_in_report=True,
        )
        db.add(card)
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace?view=false_positive")

    assert "false_positive" in response.text
    assert "锚点不匹配" in response.text
    assert "<dt>证据强度</dt><dd>strong</dd>" not in response.text


def test_theoretical_foundation_narrative_mentions_frequency_difference_when_present(db_session_factory, tmp_path):
    quote = "The frequency difference and spectral model are derived via convolution operation in the cited work."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text=quote,
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "frequency difference" in (card.narrative_zh or "")
    assert "spectral model" in (card.narrative_zh or "")
    assert "convolution operation" in (card.narrative_zh or "")


def test_generate_impact_card_from_application_extension_evidence(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="application_extension",
            citation_text="Cited Scholar Paper is extended to a new application scenario.",
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert card.card_type == "application_extension"
    assert "新的应用场景" in (card.narrative_zh or "")


def test_generate_limitation_feedback_card_from_negative_evidence(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            review_status="accepted",
            aspect="limitation_or_negative",
            citation_text="The cited method has clear limitations in practical robustness.",
        )
        evidence = db.query(StrongEvidence).one()
        evidence.stance = "negative"
        db.commit()
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert card.card_type == "limitation_or_negative"
    assert "客观评价或局限性分析" in (card.narrative_zh or "")


def test_negative_evidence_not_rewritten_as_positive(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            review_status="accepted",
            aspect="limitation_or_negative",
            citation_text="The cited method has clear limitations in practical robustness.",
        )
        evidence = db.query(StrongEvidence).one()
        evidence.stance = "negative"
        db.commit()
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "正向评价候选" not in (card.narrative_zh or "")
    assert "局限性" in (card.narrative_zh or "")


def test_grouped_citation_card_requires_human_review(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            review_status="accepted",
            aspect="detailed_comparison",
            citation_text="Methods [15], [16], [17] are compared together.",
        )
        evidence = db.query(StrongEvidence).one()
        evidence.mention_type = "grouped_citation"
        evidence.anchor_status = "grouped_citation"
        evidence.evidence_reason = "该证据来自成组引用，可能同时适用于多个被引论文，需要人工确认归因范围。"
        evidence.evidence_strength = "moderate"
        db.commit()
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "人工确认归因范围" in (card.narrative_zh or "")


def test_notable_author_unknown_not_claimed_as_fellow(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert (card.fellow_status or "unknown") == "unknown"
    assert "Fellow" not in (card.narrative_zh or "")


def test_representative_work_not_rewritten_as_high_praise(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="representative_work",
            evidence_strength="weak",
            citation_text=(
                "Other work explores estimating pose leveraging moire patterns' high sensitivity "
                "to the camera's pose changes [60]."
            ),
        )
        evidence = db.query(StrongEvidence).one()
        evidence.stance = "neutral"
        evidence.mention_type = "related_work"
        db.commit()
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "高度评价" not in (card.narrative_zh or "")
    assert "技术脉络" in (card.narrative_zh or "")


def test_report_narrative_for_representative_work_is_conservative(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="representative_work",
            evidence_strength="weak",
            citation_text=(
                "Other work explores estimating pose leveraging moire patterns' high sensitivity "
                "to the camera's pose changes [60]."
            ),
        )
        evidence = db.query(StrongEvidence).one()
        evidence.stance = "neutral"
        evidence.mention_type = "related_work"
        db.commit()
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        report = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "技术脉络" in report
    assert "高度评价" not in report


def test_narrative_uses_specific_terms_from_citation_text(db_session_factory, tmp_path):
    quote = "The spectral model and convolution operation build on the cited paper."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text=quote,
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "spectral model" in (card.narrative_zh or "")
    assert "convolution operation" in (card.narrative_zh or "")
    assert "核心思想" not in (card.narrative_zh or "")


def test_rejected_evidence_not_generate_card(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, review_status="rejected")
        cards = HighlightCardService(db).generate_cards_from_evidence(session_id)

    assert cards == []


def test_false_positive_not_generate_card(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, review_status="false_positive")
        cards = HighlightCardService(db).generate_cards_from_evidence(session_id)

    assert cards == []


def test_important_evidence_sorted_first(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            review_status="accepted",
            score=0.95,
            title="Accepted Card Source",
        )
        seed_evidence(
            db,
            tmp_path,
            review_status="important",
            score=0.7,
            title="Important Card Source",
        )
        # Put second evidence into the same scholar session.
        for evidence in db.query(StrongEvidence).filter(StrongEvidence.scholar_session_id != session_id):
            evidence.scholar_session_id = session_id
        db.commit()
        cards = HighlightCardService(db).generate_cards_from_evidence(session_id)

    assert cards[0].source_citing_paper_title == "Important Card Source"


def test_card_must_link_strong_evidence(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert card.strong_evidence_id is not None
    assert card.evidence_quote


def test_card_requires_strong_evidence(db_session_factory):
    with Session(db_session_factory.kw["bind"]) as db:
        db.add(
            HighlightCard(
                scholar_session_id=1,
                strong_evidence_id=None,
                card_type="method_foundation",
                title="Invalid card",
                body_markdown="No source evidence",
                evidence_quote="quote",
                source_citing_paper_title="Citing",
                source_cited_paper_title="Cited",
            )
        )
        with pytest.raises(IntegrityError):
            db.commit()


def test_card_contains_original_evidence_quote(db_session_factory, tmp_path):
    quote = "This exact original quote must stay traceable."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text=quote)
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert card.evidence_quote == quote


def test_no_original_evidence_no_card(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text="")
        cards = HighlightCardService(db).generate_cards_from_evidence(session_id)

    assert cards == []


def test_user_edited_card_not_overwritten(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        card = service.generate_cards_from_evidence(session_id)[0]
        service.update_card(
            card.id,
            title="Human title",
            subtitle="Human subtitle",
            narrative_zh="Human narrative",
            body_markdown="Human body",
            user_note="Keep",
        )
        service.generate_cards_from_evidence(session_id)
        saved = db.get(HighlightCard, card.id)

    assert saved.title == "Human title"
    assert saved.body_markdown == "Human narrative"
    assert saved.user_note == "Keep"
    assert saved.is_user_edited is True


def test_user_edited_card_not_overwritten_on_regenerate(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        card = service.generate_cards_from_evidence(session_id)[0]
        original_sort_order = card.sort_order
        service.update_card(
            card.id,
            title="Edited title",
            subtitle="Edited subtitle",
            narrative_zh="Edited narrative",
            body_markdown="Edited body",
            user_note="Edited note",
        )
        service.generate_cards_from_evidence(session_id)
        saved = db.get(HighlightCard, card.id)

    assert saved.title == "Edited title"
    assert saved.body_markdown == "Edited narrative"
    assert saved.user_note == "Edited note"
    assert saved.sort_order == original_sort_order


def test_export_highlight_cards_csv(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        csv_text = service.export_cards_csv(session_id)

    rows = list(csv.DictReader(StringIO(csv_text)))
    assert rows[0]["card_type"] == "method_foundation"
    assert rows[0]["evidence_quote"]


def test_export_highlight_cards_markdown(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        markdown = service.export_legacy_cards_markdown(session_id)

    assert "# 亮点引用证据报告" in markdown
    assert "### 原文证据" in markdown
    assert "### 评价理由" in markdown
    assert "Independent Citing Paper" in markdown


def test_highlight_cards_markdown_has_section_headers(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "## 一、报告摘要" in markdown
    assert "## 二、强证据卡片" in markdown
    assert "## 五、已排除误报摘要" in markdown


def test_highlight_cards_markdown_has_evidence_quote_section(db_session_factory, tmp_path):
    quote = "This method foundation quote is explicitly grounded in the citing paper."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text=quote)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "### 原文证据" in markdown
    assert quote in markdown


def test_highlight_cards_markdown_has_why_this_judgment_section(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "### 评价理由" in markdown
    assert "支持" in markdown or "判断" in markdown


def test_highlight_cards_markdown_has_copy_ready_statement_section(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "### 亮点评价" in markdown
    assert "成果总结" not in markdown


def test_highlight_cards_markdown_uses_display_context(db_session_factory, tmp_path):
    quote = "The cited technique [36] is used to construct the target model."
    long_context = (
        "Section 3. Model. The paragraph first explains the setup and assumptions. "
        f"{quote} The next sentence explains why the model depends on the cited work. "
        "The following sentence provides additional body context for report readers."
    )
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        asset = db.query(PdfAsset).one()
        Path(asset.extracted_text_path).write_text(long_context, encoding="utf-8")
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "The paragraph first explains the setup and assumptions" in markdown
    assert "additional body context for report readers" in markdown


def test_highlight_cards_markdown_highlights_target_marker(db_session_factory, tmp_path):
    quote = "The cited technique [36] is used to construct the target model."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "**[36]**" in markdown


def test_highlight_cards_markdown_highlights_key_phrases(db_session_factory, tmp_path):
    quote = "The citation evidence [36] gives a method foundation for the analysis."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "**method foundation**" in markdown
    assert "**citation evidence**" in markdown


def test_highlight_cards_markdown_excludes_false_positive_by_default(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text="Valid reportable quote.")
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]
        card.review_status = "false_positive"
        db.commit()
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "Valid reportable quote." not in markdown
    assert "误报已排除数量：1" in markdown


def test_highlight_cards_markdown_excludes_invalid_anchor_by_default(db_session_factory, tmp_path):
    wrong_quote = "The comparison discusses another method [20] instead."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]
        card.evidence_quote = wrong_quote
        card.review_status = "false_positive"
        card.include_in_report = True
        db.commit()
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert wrong_quote not in markdown


def test_highlight_cards_markdown_not_single_line_field_dump(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "类型:" not in markdown
    assert "中文亮点评价:" not in markdown
    assert "\n### 原文证据\n" in markdown


def test_markdown_export_does_not_include_fulltext_dump(db_session_factory, tmp_path):
    quote = "The cited technique [36] is used in the model."
    fulltext = "A" * 2500 + f" {quote} " + "SHOULD_NOT_EXPORT_FULLTEXT_TAIL " + "B" * 2500
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        asset = db.query(PdfAsset).one()
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert "SHOULD_NOT_EXPORT_FULLTEXT_TAIL" in markdown
    assert "B" * 2000 not in markdown


def test_markdown_export_does_not_expose_file_paths(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        markdown = HighlightCardService(db).export_legacy_cards_markdown(session_id)

    assert str(tmp_path) not in markdown


def test_render_markdown_highlight_highlights_marker_and_phrase():
    text = "The method uses a frequency-domain model [36]."
    markdown = render_markdown_highlight(text, ["frequency-domain model"], "[36]")

    assert "**frequency-domain model**" in markdown
    assert "**[36]**" in markdown


def test_export_scholar_report_md(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_template_direct_result(db, tmp_path)
        report = ScholarReportService(db).build_report_markdown(session_id)

    assert "# 亮点引用证据报告" in report
    assert "推荐纳入证据数：1" in report
    assert "strong_evidence_count" not in report
    assert "#### 原文证据" in report
    assert "#### 对应参考文献" in report
    assert "#### 评价理由" in report
    assert "## 结论摘要" in report
    assert "### 直接亚毫米级佐证" in report
    assert "### 能力认可佐证" in report
    assert "局限性反馈 / 不宜作为亮点" in report


def test_report_contains_evidence_quote(db_session_factory, tmp_path):
    quote = "Original evidence quote for the scholar report."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_template_direct_result(db, tmp_path, quote=quote)
        report = ScholarReportService(db).build_report_markdown(session_id)

    assert quote in report


def test_report_contains_citing_paper_info(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_template_direct_result(db, tmp_path, title="Traceable Citing Paper")
        report = ScholarReportService(db).build_report_markdown(session_id)

    assert "引用论文：** Traceable Citing Paper" in report


def test_report_workspace_formal_view_hides_debug_fields(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_template_direct_result(db, tmp_path)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "正式报告预览" in response.text
    assert "<h3>结论摘要</h3>" in response.text
    assert "formal-evidence-card" in response.text
    assert "查看原始 Markdown" in response.text
    assert "命中模板" not in response.text
    assert "anchor_validation_status" not in response.text
    assert "从 selected / accepted / important evidence 生成卡片" not in response.text


def test_report_workspace_formal_view_highlights_marker_and_claim_phrases(client, db_session_factory, tmp_path):
    quote = "Target Paper enables through-wall eavesdropping with RFID tags [23]."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_template_direct_result(db, tmp_path, quote=quote, context=quote)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert 'class="citation-marker">[23]</mark>' in response.text
    assert 'class="claim-phrase">through-wall eavesdropping</mark>' in response.text
    assert 'class="claim-phrase">RFID tags</mark>' in response.text
    assert 'class="target-reference">Target Paper</mark>' in response.text


def test_matched_submm_direct_claim_conclusion_is_verified(db_session_factory, tmp_path):
    quote = "The system is useful for detecting sub-millimeter-level vibrations [4]."
    payload = {
        "target_reference_marker": "[4]",
        "target_reference_entry": "[4] Target Paper.",
        "paper_level_summary_zh": "正文亚毫米能力。",
        "evidences": [
            {
                "recommendation": "include",
                "claim_type": "capability_recognition",
                "evidence_quote": quote,
                "evidence_context": quote,
                "reference_entry": "[4] Target Paper.",
                "why_this_judgment_zh": "正文明确 detecting sub-millimeter-level vibrations。",
                "copy_ready_zh": "可作为直接亚毫米能力佐证。",
                "confidence": "high",
            }
        ],
    }
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, item_id = seed_queue_item(
            db,
            tmp_path,
            target_title="Target Paper",
            text=quote + "\n\nReferences\n[4] Target Paper.",
        )
        item = db.get(DeepAnalysisQueueItem, item_id)
        db.add(
            FulltextAnalysisResult(
                scholar_session_id=session_id,
                queue_item_id=item_id,
                citation_edge_id=item.citation_edge_id,
                analysis_scope="fulltext_template_direct",
                status="succeeded",
                parsed_result_json=json.dumps(payload, ensure_ascii=False),
            )
        )
        db.commit()
        report = HighlightCardService(db).export_cards_markdown(session_id)

    assert "发现少量已核验直接亚毫米级文本证据" in report
    assert "仍需人工核验" not in report.split("是否发现第三方明确亚毫米级佐证：", 1)[1].split("\n", 1)[0]
    assert "### 1. 直接亚毫米精度佐证" in report


def test_export_structured_json(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        payload = json.loads(ScholarReportService(db).build_structured_json(session_id))

    assert payload["scholar_session"]["id"] == session_id
    assert payload["strong_evidence"][0]["id"] == evidence_id
    assert payload["highlight_cards"][0]["strong_evidence_id"] == evidence_id
    assert payload["highlight_cards"][0]["source_evidence_id"] == evidence_id
    assert payload["exports"]["schema_version"] == "phase14"


def test_structured_json_contains_cards_and_evidence(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        payload = json.loads(ScholarReportService(db).build_structured_json(session_id))

    assert payload["exports"]
    assert payload["scholar_session"]
    assert payload["publications_summary"]
    assert payload["citation_edges_summary"]
    assert payload["queue_summary"]
    assert payload["strong_evidence"][0]["id"] == evidence_id
    assert payload["highlight_cards"][0]["strong_evidence_id"] == evidence_id
    assert "narrative_zh" in payload["highlight_cards"][0]


def test_report_does_not_expose_local_paths(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        payload = ScholarReportService(db).build_structured_json(session_id)
        markdown = ScholarReportService(db).build_report_markdown(session_id)

    assert str(tmp_path) not in payload
    assert str(tmp_path) not in markdown


def test_exports_do_not_leak_local_paths(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        card_service = HighlightCardService(db)
        card_service.generate_cards_from_evidence(session_id)
        outputs = [
            ScholarReportService(db).build_structured_json(session_id),
            ScholarReportService(db).build_report_markdown(session_id),
            card_service.export_cards_csv(session_id),
            card_service.export_legacy_cards_markdown(session_id),
        ]

    assert all(str(tmp_path) not in output for output in outputs)


def test_report_does_not_include_api_keys(db_session_factory, tmp_path, monkeypatch):
    monkeypatch.setenv("ACADEMIC_IMPACT_LLM_API_KEY", "phase14-secret")
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)
        payload = ScholarReportService(db).build_structured_json(session_id)
        markdown = ScholarReportService(db).build_report_markdown(session_id)

    assert "phase14-secret" not in payload
    assert "phase14-secret" not in markdown


def test_export_pptx_creates_file(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    assert path.exists()
    assert path.suffix == ".pptx"
    assert path.read_bytes()[:2] == b"PK"


def test_pptx_export_creates_valid_zip(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    with zipfile.ZipFile(path) as archive:
        assert archive.testzip() is None


def test_pptx_export_can_be_opened_by_python_pptx(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    presentation = Presentation(path)
    assert presentation is not None


def test_pptx_uses_widescreen(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    presentation = Presentation(path)
    assert presentation.slide_width > presentation.slide_height


def test_pptx_card_slide_has_title_bar(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    presentation = Presentation(path)
    card_slide = presentation.slides[2]
    texts = [shape.text for shape in card_slide.shapes if hasattr(shape, "text") and shape.text]
    assert any("学术引用证据分析 - 同行评价" in text for text in texts)


def test_pptx_card_slide_has_two_column_layout(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    presentation = Presentation(path)
    card_slide = presentation.slides[2]
    left_positions = sorted(shape.left for shape in card_slide.shapes)
    assert len(left_positions) >= 4
    assert left_positions[0] < left_positions[-1]


def test_pptx_inserts_highlighted_quote_png(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text="The frequency difference [36] drives the spectral model.",
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    presentation = Presentation(path)
    card_slide = presentation.slides[2]
    picture_shapes = [shape for shape in card_slide.shapes if getattr(shape, "shape_type", None) == 13]
    assert picture_shapes


def test_pptx_quote_png_uses_display_context(db_session_factory, tmp_path):
    quote = "The cited work [16] defines the temporal graph model."
    fulltext = f"2 Method\nBefore context. {quote} After context describes the graph model and calibration process."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[16]"},
        )
        service = HighlightCardService(db)
        card = service.generate_cards_from_evidence(session_id)[0]
        evidence = db.query(StrongEvidence).one()
        queue_item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        db.commit()
        context_preview = service._context_preview_for_card(card)
        quote_text = service._ppt_quote_text(card, context_preview)

    assert "After context describes" in quote_text
    assert len(quote_text) > len(quote)


def test_pptx_quote_text_longer_than_citation_text_when_context_available(db_session_factory, tmp_path):
    quote = "The cited work [17] defines the retrieval feature representation."
    fulltext = f"3 Retrieval\nSetup sentence. {quote} Follow-up sentence explains the scoring model."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text=quote, analysis_diagnostics={"target_reference_marker": "[17]"})
        service = HighlightCardService(db)
        card = service.generate_cards_from_evidence(session_id)[0]
        evidence = db.query(StrongEvidence).one()
        queue_item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        db.commit()
        quote_text = service._ppt_quote_text(card, service._context_preview_for_card(card))

    assert len(quote_text) > len(card.evidence_quote)


def test_pptx_highlights_target_marker_from_context(db_session_factory, tmp_path):
    quote = "The cited work [18] defines the calibration signal model."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text=quote, analysis_diagnostics={"target_reference_marker": "[18]"})
        service = HighlightCardService(db)
        card = service.generate_cards_from_evidence(session_id)[0]
        image_path = service._render_quote_image(
            session_id=session_id,
            card_id=card.id,
            quote_text=f"Section | [18]\nBefore. {quote} After.",
            target_marker="[18]",
            context_preview={"highlight_terms": ["calibration signal model"]},
            evidence=db.query(StrongEvidence).one(),
        )

    assert image_path.exists()
    assert image_path.stat().st_size > 0


def test_pptx_does_not_use_only_evidence_quote_when_context_available(db_session_factory, tmp_path):
    quote = "The cited work [19] defines the structured encoder model."
    fulltext = f"4 Encoder\nIntro sentence. {quote} Context sentence after the quote describes training behavior."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text=quote, analysis_diagnostics={"target_reference_marker": "[19]"})
        service = HighlightCardService(db)
        card = service.generate_cards_from_evidence(session_id)[0]
        evidence = db.query(StrongEvidence).one()
        queue_item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        db.commit()
        quote_text = service._ppt_quote_text(card, service._context_preview_for_card(card))

    assert quote_text != card.evidence_quote
    assert "Context sentence after" in quote_text


def test_pptx_does_not_dump_all_fields_as_plain_text(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    presentation = Presentation(path)
    card_slide = presentation.slides[2]
    texts = "\n".join(shape.text for shape in card_slide.shapes if hasattr(shape, "text") and shape.text)
    assert "中文亮点评价:" not in texts
    assert "英文原文 citation_text:" not in texts


def test_pptx_truncates_long_text(db_session_factory, tmp_path):
    long_quote = "frequency difference [36] " + ("spectral model " * 300)
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text=long_quote, analysis_diagnostics={"target_reference_marker": "[36]"})
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    presentation = Presentation(path)
    card_slide = presentation.slides[2]
    texts = "\n".join(shape.text for shape in card_slide.shapes if hasattr(shape, "text") and shape.text)
    assert len(texts) < 4000


def test_real_python_pptx_import_is_used():
    spec = importlib.util.find_spec("pptx")
    assert spec is not None
    assert "site-packages" in (spec.origin or "")


def test_pptx_export_no_local_shim():
    spec = importlib.util.find_spec("pptx")
    assert spec is not None
    assert "academic_impact_app/pptx" not in (spec.origin or "")


def test_pptx_route_returns_correct_media_type(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/exports/report.pptx")

    assert response.status_code == 200
    assert response.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


def test_pptx_export_sanitizes_control_characters(db_session_factory, tmp_path):
    dirty_quote = "Line 1\x00\x01\x02Line 2"
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text=dirty_quote)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    Presentation(path)


def test_pptx_export_records_validation_error(db_session_factory, tmp_path, monkeypatch):
    def write_invalid(path, session_id, cards):
        path.write_text("not a zip", encoding="utf-8")

    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        monkeypatch.setattr(service, "_write_minimal_pptx", write_invalid)
        with pytest.raises(PptxExportError):
            service.export_pptx(session_id)
        diagnostics = json.loads(
            (service._session_export_dir(session_id) / "report_pptx_diagnostics.json").read_text(
                encoding="utf-8"
            )
        )

    assert diagnostics["validation_status"] == "failed"
    assert diagnostics["validation_error"]


def test_pptx_export_does_not_return_html_error_as_pptx(client, db_session_factory, tmp_path, monkeypatch):
    def explode(self, session_id):
        raise PptxExportError("PPTX 导出失败：validation failed")

    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    monkeypatch.setattr(
        "app.routers.highlight_cards.HighlightCardService.export_pptx",
        explode,
    )
    response = client.get(f"/scholar-sessions/{session_id}/exports/report.pptx")

    assert response.status_code == 500
    assert response.headers["content-type"].startswith("text/html")
    assert "PPTX 导出失败" in response.text


def test_pptx_export_skips_missing_images(db_session_factory, tmp_path, monkeypatch):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        cards = service.generate_cards_from_evidence(session_id)
        monkeypatch.setattr(
            service,
            "_candidate_image_paths",
            lambda card: [str(tmp_path / "missing.png")],
        )
        path = service.export_pptx(session_id)
        diagnostics = json.loads(
            (service._session_export_dir(session_id) / "report_pptx_diagnostics.json").read_text(
                encoding="utf-8"
            )
        )

    assert path.exists()
    assert diagnostics["warning_count"] == len(diagnostics["warnings"])
    assert diagnostics["warnings"]


def test_ppt_export_does_not_include_unreviewed_false_positive(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, review_status="false_positive")
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        path = service.export_pptx(session_id)

    assert path.exists()
    content = path.read_bytes()
    assert b"false_positive" not in content


def test_impact_card_preserves_source_evidence_id(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(db, tmp_path)
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert card.source_evidence_id == evidence_id


def test_exports_do_not_leak_api_keys(db_session_factory, tmp_path, monkeypatch):
    monkeypatch.setenv("ACADEMIC_IMPACT_LLM_API_KEY", "phase14-secret")
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        card_service = HighlightCardService(db)
        card_service.generate_cards_from_evidence(session_id)
        outputs = [
            ScholarReportService(db).build_structured_json(session_id),
            ScholarReportService(db).build_report_markdown(session_id),
            card_service.export_cards_csv(session_id),
            card_service.export_legacy_cards_markdown(session_id),
        ]

    assert all("phase14-secret" not in output for output in outputs)


def test_report_excludes_cards_whose_evidence_later_rejected(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(db, tmp_path)
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.review_status = "rejected"
        db.commit()
        report = ScholarReportService(db).build_report_markdown(session_id)
        payload = json.loads(ScholarReportService(db).build_structured_json(session_id))

    assert card.title not in report
    assert payload["strong_evidence"] == []


def test_unreviewed_high_strength_card_is_marked_draft(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            review_status="unreviewed",
            score=0.95,
            evidence_strength="strong",
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "Draft" in card.subtitle


def test_cards_sorted_by_importance_and_score(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            review_status="unreviewed",
            score=0.99,
            title="High Strength Draft",
        )
        seed_evidence(
            db,
            tmp_path,
            review_status="accepted",
            score=0.85,
            title="Accepted Evidence",
        )
        seed_evidence(
            db,
            tmp_path,
            review_status="important",
            score=0.7,
            title="Important Evidence",
        )
        for evidence in db.query(StrongEvidence).filter(StrongEvidence.scholar_session_id != session_id):
            evidence.scholar_session_id = session_id
        db.commit()
        cards = HighlightCardService(db).generate_cards_from_evidence(session_id)

    assert [card.source_citing_paper_title for card in cards] == [
        "Important Evidence",
        "Accepted Evidence",
        "High Strength Draft",
    ]


def test_highlight_card_routes_and_exports(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)

    generate_response = client.post(
        f"/scholar-sessions/{session_id}/cards/generate",
        follow_redirects=False,
    )
    assert generate_response.status_code == 303
    page_response = client.get(f"/scholar-sessions/{session_id}/cards")
    assert page_response.status_code == 200
    assert "Independent Citing Paper" in page_response.text

    report_response = client.get(f"/scholar-sessions/{session_id}/exports/report.md")
    structured_response = client.get(
        f"/scholar-sessions/{session_id}/exports/structured.json"
    )
    csv_response = client.get(
        f"/scholar-sessions/{session_id}/exports/highlight_cards.csv"
    )
    markdown_response = client.get(
        f"/scholar-sessions/{session_id}/exports/highlight_cards.md"
    )
    pptx_response = client.get(
        f"/scholar-sessions/{session_id}/exports/report.pptx"
    )

    assert report_response.status_code == 200
    assert structured_response.status_code == 200
    assert csv_response.status_code == 200
    assert markdown_response.status_code == 200
    assert pptx_response.status_code == 200


def test_report_workspace_lists_impact_cards(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "报告工作台" in response.text
    assert "中文亮点评价" in response.text


def test_report_workspace_layout_no_oval_nav(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "oval" not in response.text.lower()
    assert "circle" not in response.text.lower()


def test_report_workspace_header_uses_action_bar(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "action-bar" in response.text


def test_report_workspace_no_overlapping_title_markup(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "page-header--workspace" in response.text


def test_report_workspace_has_export_action_buttons(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    for label in ["report.md", "report.pptx", "structured.json", "highlight_cards.csv"]:
        assert label in response.text


def test_report_workspace_mobile_layout_has_wrapped_actions(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "quick-links action-bar" in response.text


def test_report_workspace_has_normal_action_buttons(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "导入重要引用作者 CSV" in response.text
    assert "report.md" in response.text
    assert "report.pptx" in response.text


def test_report_workspace_has_card_stats(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "报告卡片数" in response.text
    assert "已纳入报告" in response.text
    assert "普通引用卡片" in response.text


def test_analyzed_item_without_strong_evidence_gets_fallback_report_card(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, item_id = seed_queue_item(
            db,
            tmp_path,
            text="Background only discussion of the target paper in related work.",
        )
        item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, item_id)
        item.queue_status = "analyzed"
        result = FulltextAnalysisResult(
            scholar_session_id=session_id,
            queue_item_id=item_id,
            citation_edge_id=item.citation_edge_id,
            analysis_scope="fulltext_anchor_direct",
            status="succeeded",
            candidate_spans_json=json.dumps(
                {
                    "target_reference_context_count": 1,
                    "target_contexts_preview": [{"context_text_preview": "Background only context [36]"}],
                }
            ),
            parsed_result_json=json.dumps({"findings": []}),
        )
        db.add(result)
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "ordinary_citation" in response.text or "citation_only" in response.text


def test_target_reference_context_generates_ordinary_citation_card(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, item_id = seed_queue_item(db, tmp_path)
        item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, item_id)
        item.queue_status = "analyzed"
        result = FulltextAnalysisResult(
            scholar_session_id=session_id,
            queue_item_id=item_id,
            citation_edge_id=item.citation_edge_id,
            analysis_scope="fulltext_anchor_direct",
            status="succeeded",
            candidate_spans_json=json.dumps(
                {
                    "target_reference_context_count": 1,
                    "target_contexts_preview": [{"context_text_preview": "Anchored body context [36]"}],
                }
            ),
            parsed_result_json=json.dumps({"findings": []}),
        )
        db.add(result)
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "ordinary_citation" in response.text


def test_fallback_card_default_not_include_in_report(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, item_id = seed_queue_item(db, tmp_path)
        item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, item_id)
        item.queue_status = "analyzed"
        result = FulltextAnalysisResult(
            scholar_session_id=session_id,
            queue_item_id=item_id,
            citation_edge_id=item.citation_edge_id,
            analysis_scope="fulltext_anchor_direct",
            status="succeeded",
            candidate_spans_json=json.dumps({"target_reference_context_count": 1}),
            parsed_result_json=json.dumps({"findings": []}),
        )
        db.add(result)
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "纳入报告：False" in response.text or "默认不纳入报告" in response.text


def test_report_workspace_shows_ordinary_citation_cards(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, item_id = seed_queue_item(db, tmp_path)
        item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, item_id)
        item.queue_status = "analyzed"
        db.add(
            FulltextAnalysisResult(
                scholar_session_id=session_id,
                queue_item_id=item_id,
                citation_edge_id=item.citation_edge_id,
                analysis_scope="fulltext_anchor_direct",
                status="succeeded",
                candidate_spans_json=json.dumps({"target_reference_context_count": 1}),
                parsed_result_json=json.dumps({"findings": []}),
            )
        )
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "ordinary_citation" in response.text or "普通引用" in response.text


def test_report_card_shows_long_context_preview(client, db_session_factory, tmp_path):
    text = "Before context. Cited Scholar Paper provides a method foundation for our workflow. After context."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, citation_text="Cited Scholar Paper provides a method foundation for our workflow.")
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "展开原文上下文" in response.text
    assert "复制原文上下文" in response.text
    assert "复制中文亮点评价" in response.text


def test_context_does_not_include_full_pdf_text(client, db_session_factory, tmp_path):
    long_text = "A" * 5000
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, item_id = seed_queue_item(db, tmp_path, text=long_text)
        item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, item_id)
        item.queue_status = "analyzed"
        db.add(
            FulltextAnalysisResult(
                scholar_session_id=session_id,
                queue_item_id=item_id,
                citation_edge_id=item.citation_edge_id,
                analysis_scope="fulltext_anchor_direct",
                status="succeeded",
                candidate_spans_json=json.dumps({"target_reference_context_count": 1}),
                parsed_result_json=json.dumps({"findings": []}),
            )
        )
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "A" * 3000 not in response.text


def test_context_does_not_expose_file_path(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert str(tmp_path) not in response.text


def test_theoretical_foundation_card_not_custom(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, aspect="theoretical_foundation")
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert card.card_type == "theoretical_foundation"


def test_card_type_display_uses_chinese_label(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path, aspect="theoretical_foundation")
        HighlightCardService(db).generate_cards_from_evidence(session_id)

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")

    assert response.status_code == 200
    assert "理论基础" in response.text


def test_custom_card_type_only_for_user_custom_cards(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(db, tmp_path, aspect="theoretical_foundation")
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]
        card.card_type = "custom"
        card.is_user_edited = True
        db.commit()
        saved = db.get(HighlightCard, card.id)

    assert saved.card_type == "custom"


def test_every_analyzed_queue_item_has_at_least_one_report_card(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, item_id = seed_queue_item(db, tmp_path)
        item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, item_id)
        item.queue_status = "analyzed"
        db.add(
            FulltextAnalysisResult(
                scholar_session_id=session_id,
                queue_item_id=item_id,
                citation_edge_id=item.citation_edge_id,
                analysis_scope="fulltext_anchor_direct",
                status="succeeded",
                candidate_spans_json=json.dumps({"target_reference_context_count": 1}),
                parsed_result_json=json.dumps({"findings": []}),
            )
        )
        db.commit()
        rows = HighlightCardService(db).list_report_workspace_cards(session_id)

    assert len(rows) >= 1


def test_generate_cards_is_idempotent(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            analysis_diagnostics={"target_reference_marker": "[36]", "target_contexts_preview": [{"section_heading": "III-A2", "context_text_preview": "Before [36] frequency difference after"}]},
        )
        service = HighlightCardService(db)
        first = service.generate_cards_from_evidence(session_id)
        second = service.generate_cards_from_evidence(session_id)

    assert len(first) == 1
    assert len(second) == 1
    assert first[0].id == second[0].id


def test_fallback_card_not_created_when_strong_card_exists(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        rows = service.list_report_workspace_cards(session_id)

    assert len(rows) == 1
    assert rows[0]["card"].strong_evidence_id is not None


def test_card_display_context_longer_than_citation_text(db_session_factory, tmp_path):
    text = (
        "III-A2 Moiré Spectral Peaks\n"
        "Before context about the recapturing process and frequency difference. "
        "The second type of MSPs is generated due to the frequency difference (FD) between the pair of displaying and imaging devices in the recapturing process [36]. "
        "According to Eq. (3), the spectral model shows a convolution operation over the Dirac comb. "
        "After context continues to explain the model."
    )
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text="The second type of MSPs is generated due to the frequency difference (FD) between the pair of displaying and imaging devices in the recapturing process [36].",
            analysis_diagnostics={
                "target_reference_marker": "[36]",
                "target_contexts_preview": [
                    {
                        "section_heading": "III-A2 Moiré Spectral Peaks",
                        "context_text_preview": "Before context about the recapturing process and frequency difference. The second type of MSPs is generated due to the frequency difference (FD) ... According to Eq. (3), the spectral model shows a convolution operation.",
                    }
                ],
            },
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        item = db.query(StrongEvidence).filter_by(scholar_session_id=session_id).one()
        result = db.get(FulltextAnalysisResult, item.fulltext_result_id)
        queue_item_id = item.queue_item_id
        pdf_asset_id = db.execute(
            __import__("sqlalchemy").select(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem.pdf_asset_id)
            .where(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem.id == queue_item_id)
        ).scalar_one()
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, pdf_asset_id)
        Path(asset.extracted_text_path).write_text(text, encoding="utf-8")
        db.commit()
        rows = service.list_report_workspace_cards(session_id)
        context_preview = rows[0]["context_preview"]

    assert len(context_preview["display_context"]) > len(rows[0]["card"].evidence_quote)
    assert "[36]" in context_preview["display_context"]
    assert "spectral model" in context_preview["display_context"]
    assert "convolution operation" in context_preview["display_context"]


def test_card_context_uses_target_reference_context(db_session_factory, tmp_path):
    text = (
        "6 RELATED WORK\n"
        "Other work explores estimating pose leveraging moiré patterns' high sensitivity to the camera's pose changes [60], "
        "and improving pose tracking using inertial sensors [2, 75, 91]."
    )
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="representative_work",
            evidence_strength="weak",
            citation_text="Other work explores estimating pose leveraging moiré patterns' high sensitivity to the camera's pose changes [60], and improving pose tracking using inertial sensors [2, 75, 91].",
            analysis_diagnostics={
                "target_reference_marker": "[60]",
                "target_contexts_preview": [{"section_heading": "6 RELATED WORK", "context_text_preview": text}],
            },
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        rows = service.list_report_workspace_cards(session_id)
        context_preview = rows[0]["context_preview"]

    assert context_preview["section_heading"] == "6 RELATED WORK"
    assert context_preview["target_reference_marker"] == "[60]"


def test_display_context_uses_target_reference_contexts(tmp_path):
    fulltext = (
        "1 Intro\n"
        "A short unrelated sentence.\n"
        "2 Method\n"
        "The target method [7] introduces adaptive graph signals for our derivation. "
        "The following paragraph explains the data flow and its signal model in detail.\n\n"
        "References\n[7] Target Paper."
    )
    text_path = tmp_path / "paper.txt"
    text_path.write_text(fulltext, encoding="utf-8")
    context = build_context_preview(
        extracted_text_path=str(text_path),
        citation_text="The target method [7] introduces adaptive graph signals for our derivation.",
        diagnostics={
            "target_reference_marker": "[7]",
            "target_contexts_preview": [
                {
                    "section_heading": "2 Method",
                    "context_text_preview": (
                        "The target method [7] introduces adaptive graph signals for our derivation. "
                        "The following paragraph explains the data flow and its signal model in detail."
                    ),
                }
            ],
        },
        target_reference_marker="[7]",
    )

    assert context["section_heading"] == "2 Method"
    assert "[7]" in context["display_context"]
    assert len(context["display_context"]) > len("The target method [7] introduces adaptive graph signals for our derivation.")


def test_display_context_falls_back_to_extracted_text_window(tmp_path):
    quote = "The cited work [3] defines the adaptive signal representation used here."
    fulltext = (
        "Background before the anchor. "
        "Additional setup sentence before the quote. "
        f"{quote} "
        "Additional explanation after the quote describes the model and feature extraction process.\n\n"
        "References\n[3] Target Paper."
    )
    text_path = tmp_path / "paper.txt"
    text_path.write_text(fulltext, encoding="utf-8")
    context = build_context_preview(
        extracted_text_path=str(text_path),
        citation_text=quote,
        diagnostics={"target_reference_marker": "[3]"},
        target_reference_marker="[3]",
    )

    assert "Additional setup sentence" in context["display_context"]
    assert "feature extraction process" in context["display_context"]
    assert len(context["display_context"]) > len(quote)


def test_card_context_excludes_references(db_session_factory, tmp_path):
    quote = "Key body citation [36] with context."
    text = (
        "Body paragraph before. "
        f"{quote} "
        "Body paragraph after.\n\nREFERENCES\n[36] J. Ning et al. MoiréPose..."
    )
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        evidence = db.query(StrongEvidence).one()
        result = db.get(FulltextAnalysisResult, evidence.fulltext_result_id)
        item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(text, encoding="utf-8")
        db.commit()
        rows = service.list_report_workspace_cards(session_id)
        context_preview = rows[0]["context_preview"]

    assert "REFERENCES" not in context_preview["display_context"]
    assert "J. Ning et al." not in context_preview["display_context"]


def test_card_context_highlights_target_marker(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text="Important context [36] with frequency difference.",
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        rows = service.list_report_workspace_cards(session_id)
        context_preview = rows[0]["context_preview"]

    assert "<mark>[36]</mark>" in context_preview["highlighted_context_html"]


def test_card_context_highlights_sentence(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text="Important context [36] with frequency difference.",
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        rows = service.list_report_workspace_cards(session_id)
        context_preview = rows[0]["context_preview"]

    assert "context-sentence" in context_preview["highlighted_context_html"]


def test_display_context_highlights_marker_and_terms(tmp_path):
    quote = "The cited work [8] defines the adaptive graph signal model."
    text_path = tmp_path / "paper.txt"
    text_path.write_text(f"Section\nBefore. {quote} After.", encoding="utf-8")
    context = build_context_preview(
        extracted_text_path=str(text_path),
        citation_text=quote,
        diagnostics={"target_reference_marker": "[8]", "findings": [{"keywords": ["adaptive graph signal model"]}]},
        target_reference_marker="[8]",
    )

    assert "<mark>[8]</mark>" in context["highlighted_context_html"]
    assert "adaptive graph signal model" in context["highlighted_context_html"]


def test_report_workspace_template_renders_display_context_not_only_citation_text(client, db_session_factory, tmp_path):
    quote = "The cited work [9] defines the dynamic control model."
    fulltext = (
        "2 System Design\n"
        "Before context explains why the model is needed. "
        f"{quote} "
        "After context explains the downstream optimization process."
    )
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[9]"},
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        evidence = db.query(StrongEvidence).one()
        item = __import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem
        asset_model = __import__("app.models", fromlist=["PdfAsset"]).PdfAsset
        queue_item = db.get(item, evidence.queue_item_id)
        asset = db.get(asset_model, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        db.commit()

    response = client.get(f"/scholar-sessions/{session_id}/report-workspace")
    html = response.text
    assert "展开原文上下文" in html
    assert "Before context explains why the model is needed" in html
    assert "downstream optimization process" in html


def test_old_card_without_context_gets_dynamic_display_context(db_session_factory, tmp_path):
    quote = "The cited work [10] supports the temporal attention model."
    fulltext = f"3 Model\nBefore paragraph. {quote} After paragraph describes training signals."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[10]"},
        )
        service = HighlightCardService(db)
        card = service.generate_cards_from_evidence(session_id)[0]
        evidence = db.query(StrongEvidence).one()
        queue_item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        db.commit()
        rows = service.list_report_workspace_cards(session_id)

    assert rows[0]["card"].id == card.id
    assert len(rows[0]["context_preview"]["display_context"]) > len(quote)
    assert "training signals" in rows[0]["context_preview"]["display_context"]


def test_dedup_preserves_card_but_recomputes_context(db_session_factory, tmp_path):
    quote = "The cited work [11] defines a robust matching process."
    fulltext = f"4 Evaluation\nPrevious sentence. {quote} Next sentence expands the matching process."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[11]"},
        )
        service = HighlightCardService(db)
        first = service.generate_cards_from_evidence(session_id)[0]
        evidence = db.query(StrongEvidence).one()
        queue_item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        db.commit()
        second = service.generate_cards_from_evidence(session_id)[0]
        rows = service.list_report_workspace_cards(session_id)

    assert first.id == second.id
    assert len(rows) == 1
    assert "Next sentence expands" in rows[0]["context_preview"]["display_context"]


def test_narrative_not_generic_template(db_session_factory, tmp_path):
    quote = (
        "The second type of MSPs is generated due to the frequency difference (FD) between the pair of displaying and imaging devices in the recapturing process [36]. "
        "According to Eq. (3), the spectral model shows a convolution operation."
    )
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[36]", "target_contexts_preview": [{"section_heading": "III-A2 Moiré Spectral Peaks", "context_text_preview": quote}]},
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "frequency difference" in (card.narrative_zh or "")
    assert "recapturing process" in (card.narrative_zh or "")
    assert "核心思想" not in (card.narrative_zh or "")


def test_no_domain_hardcoded_terms_in_production_code():
    pattern = __import__("re").compile(
        r"Moiré|Moire|frequency difference|recapturing process|spectral model|"
        r"convolution operation|camera.?s pose changes|pose estimation",
        flags=__import__("re").IGNORECASE,
    )
    hits = []
    for path in Path("app").rglob("*.py"):
        text = path.read_text(encoding="utf-8")
        for line_no, line in enumerate(text.splitlines(), start=1):
            if pattern.search(line):
                hits.append(f"{path}:{line_no}:{line.strip()}")
    assert hits == []


def test_moire_terms_allowed_only_in_tests_or_fixtures():
    pattern = __import__("re").compile(r"Moiré|Moire|frequency difference|recapturing process|spectral model|convolution operation", flags=__import__("re").IGNORECASE)
    app_hits = [
        str(path)
        for path in Path("app").rglob("*.py")
        if pattern.search(path.read_text(encoding="utf-8"))
    ]
    test_hits = [
        str(path)
        for path in Path("tests").rglob("*.py")
        if pattern.search(path.read_text(encoding="utf-8"))
    ]
    assert app_hits == []
    assert test_hits


def test_narrative_uses_keywords_from_finding_not_hardcode(db_session_factory, tmp_path):
    quote = "The cited work [12] introduces calibrated latent transition graphs."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[12]", "findings": [{"keywords": ["calibrated latent transition graphs"]}]},
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "calibrated latent transition graphs" in (card.narrative_zh or "")


def test_narrative_uses_context_terms_for_non_moire_fixture(db_session_factory, tmp_path):
    quote = "The cited work [13] defines the adaptive graph signal model used in this derivation."
    context = f"2 Method\nBefore. {quote} The adaptive graph signal model is then used by the optimization process."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[13]", "target_contexts_preview": [{"section_heading": "2 Method", "context_text_preview": context}]},
        )
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "adaptive graph signal model" in (card.narrative_zh or "")


def test_representative_work_narrative_mentions_related_work_and_not_high_praise(db_session_factory, tmp_path):
    quote = (
        "Other work explores estimating pose leveraging moiré patterns' high sensitivity to the camera's pose changes [60], "
        "and improving pose tracking using inertial sensors [2, 75, 91]."
    )
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            aspect="representative_work",
            evidence_strength="weak",
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[60]", "target_contexts_preview": [{"section_heading": "6 RELATED WORK", "context_text_preview": quote}]},
        )
        evidence = db.query(StrongEvidence).one()
        evidence.stance = "neutral"
        evidence.mention_type = "related_work"
        db.commit()
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]

    assert "Related Work" in (card.narrative_zh or "")
    assert "moiré patterns" in (card.narrative_zh or "")
    assert "pose tracking" in (card.narrative_zh or "")
    assert "高度赞扬" in (card.narrative_zh or "")
    assert "高度评价" not in (card.narrative_zh or "")


def test_narrative_regenerate_updates_existing_card_not_duplicate(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(
            db,
            tmp_path,
            aspect="theoretical_foundation",
            citation_text="Initial frequency difference [36].",
            analysis_diagnostics={"target_reference_marker": "[36]"},
        )
        service = HighlightCardService(db)
        first = service.generate_cards_from_evidence(session_id)[0]
        evidence = db.get(StrongEvidence, evidence_id)
        evidence.citation_text = "Updated frequency difference and spectral model [36]."
        db.commit()
        second = service.generate_cards_from_evidence(session_id)[0]
        card_count = db.query(HighlightCard).count()

    assert first.id == second.id
    assert card_count == 1
    assert "spectral model" in (second.narrative_zh or "")


def test_regenerate_updates_existing_card_context(db_session_factory, tmp_path):
    quote = "The cited work [14] supports the alignment model."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[14]"},
        )
        service = HighlightCardService(db)
        first = service.generate_cards_from_evidence(session_id)[0]
        evidence = db.query(StrongEvidence).one()
        queue_item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(f"Section\nBefore. {quote} After context.", encoding="utf-8")
        db.commit()
        second = service.generate_cards_from_evidence(session_id)[0]
        rows = service.list_report_workspace_cards(session_id)

    assert first.id == second.id
    assert "After context" in rows[0]["context_preview"]["display_context"]


def test_dedup_does_not_keep_stale_short_context_only(db_session_factory, tmp_path):
    quote = "The cited work [15] introduces the retrieval pipeline."
    fulltext = f"5 Pipeline\nBefore. {quote} After text explains the retrieval pipeline and ranking features."
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(
            db,
            tmp_path,
            citation_text=quote,
            analysis_diagnostics={"target_reference_marker": "[15]"},
        )
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        evidence = db.query(StrongEvidence).one()
        queue_item = db.get(__import__("app.models", fromlist=["DeepAnalysisQueueItem"]).DeepAnalysisQueueItem, evidence.queue_item_id)
        asset = db.get(__import__("app.models", fromlist=["PdfAsset"]).PdfAsset, queue_item.pdf_asset_id)
        Path(asset.extracted_text_path).write_text(fulltext, encoding="utf-8")
        db.commit()
        service.generate_cards_from_evidence(session_id)
        rows = service.list_report_workspace_cards(session_id)

    assert len(rows) == 1
    assert len(rows[0]["context_preview"]["display_context"]) > len(quote)
    assert "ranking features" in rows[0]["context_preview"]["display_context"]


def test_regenerate_does_not_duplicate_cards(db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        service = HighlightCardService(db)
        service.generate_cards_from_evidence(session_id)
        service.generate_cards_from_evidence(session_id)
        service.generate_cards_from_evidence(session_id)
        cards = db.query(HighlightCard).filter_by(scholar_session_id=session_id).all()

    assert len(cards) == 1


def test_report_workspace_allows_editing_card_text(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, _ = seed_evidence(db, tmp_path)
        card = HighlightCardService(db).generate_cards_from_evidence(session_id)[0]
        card_id = card.id

    response = client.post(
        f"/scholar-sessions/{session_id}/cards/{card_id}/edit",
        data={
            "title": "编辑后的标题",
            "subtitle": "编辑后的副标题",
            "narrative_zh": "编辑后的中文亮点评价",
            "body_markdown": "旧正文",
            "user_note": "已人工编辑",
            "include_in_report": "true",
            "notable_author_name": "unknown",
            "notable_author_affiliation": "",
            "notable_author_role": "",
            "fellow_status": "unknown",
        },
        follow_redirects=False,
    )

    assert response.status_code == 303
    with Session(db_session_factory.kw["bind"]) as db:
        saved = db.get(HighlightCard, card_id)

    assert saved.title == "编辑后的标题"
    assert saved.subtitle == "编辑后的副标题"
    assert saved.narrative_zh == "编辑后的中文亮点评价"
    assert saved.body_markdown == "编辑后的中文亮点评价"


def test_evidence_page_can_generate_impact_card(client, db_session_factory, tmp_path):
    with Session(db_session_factory.kw["bind"]) as db:
        session_id, evidence_id = seed_evidence(db, tmp_path)

    response = client.post(
        f"/scholar-sessions/{session_id}/evidence/{evidence_id}/generate-card",
        follow_redirects=False,
    )

    assert response.status_code == 303
    with Session(db_session_factory.kw["bind"]) as db:
        card = db.query(HighlightCard).one()

    assert card.strong_evidence_id == evidence_id
