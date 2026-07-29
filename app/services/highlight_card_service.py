"""Service for generating and exporting scholar highlight cards."""

import csv
from dataclasses import dataclass
import html
import json
from io import StringIO
import os
from pathlib import Path
import re
import shutil
from types import SimpleNamespace
import tempfile
import unicodedata
import zipfile
from typing import List, Optional, Tuple

from fastapi import Depends
from sqlalchemy.orm import Session
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.analysis.card_builder import build_card_values, card_type_for_evidence, generate_impact_narrative
from app.analysis.citation_anchor import reference_entries_by_marker
from app.analysis.target_anchor_validation import validate_citation_target_anchor
from app.analysis.template_direct_postprocess import postprocess_template_direct_payload
from app.core.config import settings
from app.db.session import get_db
from app.models import (
    AnalysisTemplate,
    CitationAuthorAnnotation,
    DeepAnalysisQueueItem,
    FulltextAnalysisResult,
    HighlightCard,
    NotableAuthor,
    PdfAsset,
    Publication,
    ScholarAnalysisSession,
    StrongEvidence,
)
from app.repositories.highlight_card_repo import HighlightCardRepository
from app.services.context_service import build_context_preview
from app.services.template_service import TemplateService


class PptxExportError(ValueError):
    pass


@dataclass
class ReportCardViewModel:
    card_index: int
    card_type_label_zh: str
    title_zh: str
    citing_paper_title: str
    citing_venue: str
    citing_year: str
    cited_paper_title: str
    evidence_type_label_zh: str
    evidence_strength: str
    stance: str
    notable_author_name: str
    honor_category: str
    citation_text: str
    display_context: str
    highlighted_citation_markdown: str
    highlighted_context_markdown: str
    key_phrases: List[str]
    target_reference_marker: str
    why_this_judgment: str
    limitation_zh: str
    copy_ready_statement: str
    report_recommendation: str
    risk_note: str
    confidence_level: str
    anchor_validation_status: str
    anchor_validation_reason: str
    include_in_report: bool
    review_status: str
    matched_template_names: List[str]
    template_match_reason: str
    template_satisfied: Optional[bool]
    template_failure_reason: str


def render_markdown_highlight(
    text: str,
    highlights: Optional[List[str]] = None,
    target_marker: Optional[str] = None,
) -> str:
    """Render safe, non-overlapping Markdown bold highlights for report exports."""
    source = str(text or "")
    if not source:
        return ""
    terms: List[str] = []
    if target_marker:
        terms.append(str(target_marker))
    for term in highlights or []:
        value = str(term or "").strip()
        if value and (len(value) >= 3 or value == target_marker):
            terms.append(value)
    deduped_terms = []
    seen = set()
    for term in sorted(terms, key=len, reverse=True):
        key = term.casefold()
        if key not in seen:
            deduped_terms.append(term)
            seen.add(key)
    spans: List[tuple[int, int]] = []
    for term in deduped_terms[:16]:
        for match in re.finditer(re.escape(term), source, flags=re.IGNORECASE):
            start, end = match.span()
            if start == end:
                continue
            if any(not (end <= existing_start or start >= existing_end) for existing_start, existing_end in spans):
                continue
            spans.append((start, end))
    if not spans:
        return source
    parts: List[str] = []
    cursor = 0
    for start, end in sorted(spans):
        parts.append(source[cursor:start])
        parts.append(f"**{source[start:end]}**")
        cursor = end
    parts.append(source[cursor:])
    return "".join(parts)


class HighlightCardService:
    def __init__(self, db: Session) -> None:
        self.repository = HighlightCardRepository(db)
        self.db = db

    def generate_cards_from_evidence(self, session_id: int) -> List[HighlightCard]:
        rows = self._dedupe_evidence_rows(self.repository.list_eligible_evidence(session_id))
        existing_cards = self.repository.list_cards(session_id)
        cards: List[HighlightCard] = []
        seen_card_ids = set()
        for sort_order, (evidence, item) in enumerate(rows, start=1):
            validation = self._anchor_validation_for_evidence(evidence, item)
            if validation.anchor_validation_status != "unknown" and not validation.is_valid:
                existing = self.repository.find_by_evidence_id(evidence.id)
                self._mark_evidence_and_card_false_positive(evidence, existing, validation)
                continue
            notable_author = self._find_notable_author_for_item(item.id, item.citing_authors_json)
            context_preview = self._context_preview_for_evidence(evidence)
            values = build_card_values(
                evidence=evidence,
                item=item,
                sort_order=sort_order,
                context_preview=context_preview,
                notable_author=notable_author,
            )
            values.update(self._template_card_values(evidence))
            matched_template_type = self._best_matched_template_type(evidence.id)
            if (
                matched_template_type
                and matched_template_type != "custom"
                and not (
                    (evidence.aspect or "") == "representative_work"
                    and matched_template_type != "representative_work"
                )
            ):
                values["card_type"] = matched_template_type
            existing = self._select_keeper_card(existing_cards, evidence, values["card_type"])
            if existing is None:
                card = self.repository.create_card(**values)
            else:
                if not existing.is_user_edited:
                    for key, value in values.items():
                        setattr(existing, key, value)
                else:
                    existing.sort_order = sort_order
                card = existing
                seen_card_ids.add(card.id)
            cards.append(card)
        self._deduplicate_existing_cards(session_id)
        self.db.commit()
        return self.list_cards(session_id)

    def generate_card_from_evidence(
        self,
        session_id: int,
        evidence_id: int,
    ) -> HighlightCard:
        self.generate_cards_from_evidence(session_id)
        card = self.repository.find_by_evidence_id(evidence_id)
        if card is None or self._is_false_positive_card(card):
            raise ValueError(f"No reportable card could be generated for evidence {evidence_id}")
        return card

    def list_cards(
        self,
        session_id: int,
        card_type: Optional[str] = None,
        view: str = "all",
    ) -> List[HighlightCard]:
        self._mark_mismatched_cards_false_positive(session_id)
        cards = self.repository.list_cards(session_id)
        if self._is_false_positive_view(view):
            cards = [card for card in cards if self._is_false_positive_card(card)]
        else:
            cards = [
                card
                for card in cards
                if card.review_status not in {"false_positive", "rejected"}
                and not self._is_invalid_anchor_card(card)
            ]
        if card_type:
            cards = [card for card in cards if card.card_type == card_type]
        return cards

    def list_report_workspace_cards(
        self,
        session_id: int,
        card_type: Optional[str] = None,
        view: str = "all",
    ) -> List[dict]:
        cards = self.list_cards(session_id, card_type=card_type, view=view)
        if view.startswith("template:"):
            template_id = view.split(":", 1)[1]
            cards = [
                card for card in cards
                if template_id in {str(value) for value in self._load_json_list(card.matched_template_ids_json)}
            ]
        elif view == "template_satisfied":
            cards = [card for card in cards if bool(card.template_satisfied)]
        elif view == "template_unsatisfied":
            cards = [card for card in cards if card.template_satisfied is False]
        rows = []
        covered_queue_item_ids = set()
        for card in cards:
            evidence = self.db.get(StrongEvidence, card.strong_evidence_id) if card.strong_evidence_id else None
            if evidence and evidence.queue_item_id:
                covered_queue_item_ids.add(evidence.queue_item_id)
            narrative_meta = self._narrative_meta_for_card(card, evidence, context_preview=self._context_preview_for_card(card))
            rows.append(
                {
                    "card": card,
                    "card_kind": (
                        "false_positive"
                        if self._is_false_positive_card(card)
                        else "strong"
                        if card.evidence_strength in {"strong", "moderate"}
                        else "ordinary"
                    ),
                    "grouped_warning": card.aspect == "limitation_or_negative"
                    or card.card_type == "limitation_or_negative"
                    or card.review_status in {"needs_discussion", "unreviewed"}
                    or "成组引用" in (card.narrative_zh or ""),
                    "context_preview": self._context_preview_for_card(card),
                    "narrative_meta": narrative_meta,
                }
            )
        for row in self._fallback_workspace_rows(session_id, covered_queue_item_ids, view=view):
            rows.append(row)
        if card_type:
            rows = [row for row in rows if row["card"].card_type == card_type]
        return rows

    def report_workspace_stats(self, session_id: int) -> dict:
        rows = self.list_report_workspace_cards(session_id)
        return {
            "report_card_count": len(rows),
            "include_in_report_count": sum(1 for row in rows if row["card"].include_in_report),
            "strong_evidence_card_count": sum(1 for row in rows if row["card_kind"] == "strong"),
            "ordinary_citation_card_count": sum(
                1
                for row in rows
                if row["card"].card_type in {"ordinary_citation", "background_reference", "citation_only", "weak_mention", "representative_work"}
            ),
            "important_author_count": sum(
                1 for row in rows if (row["card"].fellow_status or "unknown") != "unknown"
            ),
            "needs_review_count": sum(
                1 for row in rows if row["card"].review_status in {"unreviewed", "needs_discussion"}
            ),
        }

    def update_card(
        self,
        card_id: int,
        *,
        title: str,
        subtitle: Optional[str] = None,
        narrative_zh: Optional[str] = None,
        body_markdown: str,
        user_note: str,
        include_in_report: bool = True,
        notable_author_name: Optional[str] = None,
        notable_author_affiliation: Optional[str] = None,
        notable_author_role: Optional[str] = None,
        fellow_status: Optional[str] = None,
    ) -> HighlightCard:
        card = self.repository.get_card(card_id)
        if card is None:
            raise ValueError(f"HighlightCard {card_id} was not found")
        card.title = title
        if subtitle is not None:
            card.subtitle = subtitle
        if narrative_zh is not None:
            card.narrative_zh = narrative_zh
            card.body_markdown = narrative_zh
        else:
            card.body_markdown = body_markdown
        card.user_note = user_note
        card.include_in_report = include_in_report
        if notable_author_name is not None:
            card.notable_author_name = notable_author_name
        if notable_author_affiliation is not None:
            card.notable_author_affiliation = notable_author_affiliation
        if notable_author_role is not None:
            card.notable_author_role = notable_author_role
        if fellow_status is not None:
            card.fellow_status = fellow_status
        card.is_user_edited = True
        self.db.commit()
        self.db.refresh(card)
        return card

    def reorder_cards(self, session_id: int, card_ids: List[int]) -> List[HighlightCard]:
        cards_by_id = {card.id: card for card in self.repository.list_cards(session_id)}
        for sort_order, card_id in enumerate(card_ids, start=1):
            card = cards_by_id.get(card_id)
            if card is not None:
                card.sort_order = sort_order
        self.db.commit()
        return self.list_cards(session_id)

    def export_cards_csv(self, session_id: int) -> str:
        output = StringIO()
        fieldnames = [
            "id",
            "card_type",
            "title",
            "subtitle",
            "source_citing_paper_title",
            "source_cited_paper_title",
            "aspect",
            "stance",
            "evidence_strength",
            "score",
            "evidence_quote",
        ]
        writer = csv.DictWriter(output, fieldnames=fieldnames)
        writer.writeheader()
        for card in self.list_cards(session_id):
            writer.writerow({field: getattr(card, field) for field in fieldnames})
        return output.getvalue()

    def export_cards_markdown(self, session_id: int, *, include_all: bool = False) -> str:
        return self._export_template_direct_markdown(session_id, include_all=include_all)

    def build_formal_report_view(self, session_id: int, *, include_all: bool = False) -> dict:
        model = self._build_template_direct_report_model(session_id, include_all=include_all)
        if not model["has_results"]:
            return {
                "has_results": False,
                "markdown": self._render_template_direct_markdown(model),
                "sections": [],
                "summary": model["summary"],
                "conclusion": model["conclusion"],
            }
        return {
            "has_results": True,
            "markdown": self._render_template_direct_markdown(model),
            "sections": model["sections"],
            "summary": model["summary"],
            "conclusion": model["conclusion"],
        }

    def build_formal_evidence_view(self, evidence_rows: List[dict]) -> dict:
        """Build the shared formal-card view model for the StrongEvidence page."""
        sections = {
            "include": {"title": "推荐纳入", "items": []},
            "submm_review": {"title": "直接亚毫米级佐证候选", "items": []},
            "review": {"title": "候选复核", "items": []},
            "limitation": {"title": "局限性 / 不宜作为亮点", "items": []},
            "exclude": {"title": "不纳入", "items": []},
        }
        for index, row in enumerate(evidence_rows, start=1):
            formal_item = self._strong_evidence_formal_item(row, index=index)
            sections[formal_item["section_key"]]["items"].append(formal_item)
        return {
            "sections": [section for section in sections.values() if section["items"]],
            "item_count": len(evidence_rows),
        }

    def _strong_evidence_formal_item(self, row: dict, *, index: int) -> dict:
        evidence = row["evidence"]
        item = row["item"]
        context_preview = row.get("context_preview") or {}
        judgment = row.get("judgment_basis") or {}
        marker = (
            judgment.get("target_reference_marker")
            or context_preview.get("target_reference_marker")
            or ""
        )
        reference_entry = self._reference_entry_for_strong_evidence(
            evidence=evidence,
            item=item,
            marker=marker,
        )
        quote = evidence.citation_text or ""
        context = (
            context_preview.get("display_context")
            or judgment.get("evidence_context")
            or quote
        )
        extra_highlights = judgment.get("key_phrases") or self._load_json_list(evidence.highlight_keywords_json)
        recommendation_label, section_key = self._formal_evidence_recommendation(evidence, judgment)
        card_type = card_type_for_evidence(evidence)
        claim_label = (
            judgment.get("judgment_label")
            or self._template_direct_claim_label(evidence.aspect or card_type)
        )
        return {
            "index": index,
            "evidence_id": evidence.id,
            "claim_label": claim_label,
            "citing_paper_title": item.citing_paper_title,
            "cited_paper_title": item.cited_paper_title,
            "recommendation_label": recommendation_label,
            "confidence": judgment.get("confidence_level") or judgment.get("confidence") or "medium",
            "citation_text": quote,
            "quote_html": self._highlight_template_direct_html(
                quote,
                marker=marker,
                reference_entry=reference_entry,
                item=item,
                extra_highlights=extra_highlights,
            ),
            "reference_entry_html": self._highlight_template_direct_html(
                reference_entry,
                marker=marker,
                reference_entry=reference_entry,
                item=item,
                extra_highlights=[],
            ),
            "context_html": self._highlight_template_direct_html(
                context,
                marker=marker,
                reference_entry=reference_entry,
                item=item,
                extra_highlights=extra_highlights,
            ),
            "evaluation_zh": judgment.get("narrative_zh") or judgment.get("evidence_claim_zh") or evidence.evidence_reason,
            "why_this_judgment_zh": judgment.get("judgment_basis_zh") or judgment.get("why_this_judgment") or evidence.evidence_reason,
            "copy_ready_zh": judgment.get("copy_ready_statement_zh") or judgment.get("copy_ready_statement") or "",
            "section_key": section_key,
        }

    def _formal_evidence_recommendation(self, evidence, judgment: dict) -> Tuple[str, str]:
        if evidence.review_status in {"false_positive", "rejected"}:
            return "不纳入", "exclude"
        if evidence.aspect == "limitation_or_negative" or evidence.stance in {"negative", "mixed"}:
            return "局限性反馈", "limitation"
        recommendation = str(judgment.get("report_recommendation") or "")
        if recommendation == "推荐纳入":
            return "推荐纳入", "include"
        if "亚毫米" in str(judgment.get("judgment_label") or ""):
            return "候选复核", "submm_review"
        if recommendation in {"不建议纳入", "不纳入"}:
            return "不纳入", "exclude"
        return "候选复核", "review"

    def _reference_entry_for_strong_evidence(self, *, evidence, item, marker: str) -> str:
        marker_number = str(marker or "").strip().strip("[]")
        if marker_number:
            entry = self._reference_entries_for_item(item).get(marker_number)
            if entry:
                return entry if entry.lstrip().startswith(f"[{marker_number}]") else f"[{marker_number}] {entry}"
        result = self.db.get(FulltextAnalysisResult, evidence.fulltext_result_id)
        if result is not None:
            parsed = self._load_json(result.parsed_result_json)
            entry = str(parsed.get("target_reference_entry") or "").strip()
            if entry:
                return entry
        marker_text = f"[{marker_number}]" if marker_number else "目标引用"
        return f"{marker_text} 未解析到引用论文原文 References 条目"

    def has_template_direct_results(self, session_id: int) -> bool:
        return (
            self.db.query(FulltextAnalysisResult.id)
            .filter(
                FulltextAnalysisResult.scholar_session_id == session_id,
                FulltextAnalysisResult.analysis_scope == "fulltext_template_direct",
                FulltextAnalysisResult.status == "succeeded",
            )
            .first()
            is not None
        )

    def export_legacy_cards_markdown(self, session_id: int, *, include_all: bool = False) -> str:
        """Export the pre-fulltext_template_direct card report for debug tooling only.

        Formal Markdown reports intentionally do not call this path because it
        includes legacy card/template diagnostics that are not suitable for the
        user-facing report.
        """
        rows = self.list_report_workspace_cards(
            session_id,
            view="debug" if include_all else "all",
        )
        view_models = [
            self._report_card_view_model(index, row)
            for index, row in enumerate(rows, start=1)
            if self._should_export_markdown_row(row, include_all=include_all)
        ]
        excluded_rows = self.list_report_workspace_cards(session_id, view="false_positive")
        session = self.db.get(ScholarAnalysisSession, session_id)
        lines = [
            "# 亮点引用证据报告",
            "",
            "## 一、报告摘要",
            "",
            f"- 学者会话：{session.display_name if session else session_id}",
            f"- 目标论文数量：{len({vm.cited_paper_title for vm in view_models if vm.cited_paper_title})}",
            f"- 报告卡片数量：{len(view_models)}",
            f"- 强证据数量：{sum(1 for vm in view_models if vm.evidence_strength in {'strong', 'moderate'})}",
            f"- 普通引用数量：{sum(1 for vm in view_models if vm.evidence_strength not in {'strong', 'moderate'})}",
            f"- 需要复核数量：{sum(1 for vm in view_models if vm.review_status in {'unreviewed', 'needs_discussion'} or vm.risk_note)}",
            f"- 误报已排除数量：{len(excluded_rows)}",
            "",
        ]
        if not view_models:
            lines.extend(
                [
                    "暂无可导出的报告卡片。",
                    "",
                    "## 五、已排除误报摘要",
                    "",
                    self._excluded_false_positive_summary(excluded_rows),
                    "",
                ]
            )
            return "\n".join(lines).rstrip() + "\n"

        negative_cards = [
            vm for vm in view_models
            if vm.card_type_label_zh == "局限性反馈" or vm.stance in {"negative", "mixed"}
        ]
        negative_keys = {(vm.title_zh, vm.citation_text, vm.card_type_label_zh) for vm in negative_cards}
        ordinary_cards = [
            vm for vm in view_models
            if (vm.title_zh, vm.citation_text, vm.card_type_label_zh) not in negative_keys
            and (
                vm.card_type_label_zh in {"普通引用", "背景引用", "代表性相关工作"}
                or vm.evidence_strength not in {"strong", "moderate"}
            )
            and vm.report_recommendation != "推荐纳入"
        ]
        ordinary_keys = {(vm.title_zh, vm.citation_text, vm.card_type_label_zh) for vm in ordinary_cards}
        high_confidence_cards = [
            vm for vm in view_models
            if (vm.title_zh, vm.citation_text, vm.card_type_label_zh) not in negative_keys
            and (vm.title_zh, vm.citation_text, vm.card_type_label_zh) not in ordinary_keys
            and (
                vm.report_recommendation == "推荐纳入"
                or (vm.evidence_strength in {"strong", "moderate"} and "成组引用" not in vm.risk_note)
            )
        ]
        high_keys = {(vm.title_zh, vm.citation_text, vm.card_type_label_zh) for vm in high_confidence_cards}
        candidate_cards = [
            vm for vm in view_models
            if (vm.title_zh, vm.citation_text, vm.card_type_label_zh) not in negative_keys
            and (vm.title_zh, vm.citation_text, vm.card_type_label_zh) not in ordinary_keys
            and (vm.title_zh, vm.citation_text, vm.card_type_label_zh) not in high_keys
        ]
        self._append_markdown_card_section(lines, "## 二、强证据卡片 / 高可信第三方佐证", high_confidence_cards)
        self._append_markdown_card_section(lines, "## 三、候选佐证，需要人工复核", candidate_cards)
        self._append_markdown_card_section(lines, "## 四、局限性/负面反馈", negative_cards)
        self._append_markdown_card_section(lines, "## 五、普通相关工作引用", ordinary_cards)
        excluded_heading = "## 六、已排除误报及原因"
        lines.extend(
            [
                excluded_heading,
                "",
                "<!-- ## 五、已排除误报摘要 -->",
                "",
                self._excluded_false_positive_summary(excluded_rows),
                "",
            ]
        )
        return "\n".join(lines).rstrip() + "\n"

    def _export_template_direct_markdown(self, session_id: int, *, include_all: bool = False) -> str:
        return self._render_template_direct_markdown(
            self._build_template_direct_report_model(session_id, include_all=include_all)
        )

    def _build_template_direct_report_model(self, session_id: int, *, include_all: bool = False) -> dict:
        session = self.db.get(ScholarAnalysisSession, session_id)
        results = (
            self.db.query(FulltextAnalysisResult)
            .filter(
                FulltextAnalysisResult.scholar_session_id == session_id,
                FulltextAnalysisResult.analysis_scope == "fulltext_template_direct",
                FulltextAnalysisResult.status == "succeeded",
            )
            .order_by(FulltextAnalysisResult.created_at.desc(), FulltextAnalysisResult.id.desc())
            .all()
        )
        payloads = []
        for result in results:
            payload = self._load_json(result.parsed_result_json)
            if isinstance(payload.get("evidences"), list):
                item = self.db.get(DeepAnalysisQueueItem, result.queue_item_id) if result.queue_item_id else None
                payloads.append((result, item, self._normalized_template_direct_payload(result, item, payload)))
        if not payloads:
            return {
                "has_results": False,
                "session_label": session.display_name if session else session_id,
                "summary": {
                    "session_label": session.display_name if session else session_id,
                    "include_count": 0,
                    "review_count": 0,
                    "submm_review_count": 0,
                    "limitation_count": 0,
                    "exclude_count": 0,
                },
                "conclusion": {
                    "submm": "尚未运行 fulltext_template_direct 分析。",
                    "first": "尚未运行 fulltext_template_direct 分析。",
                    "capability": "尚未运行 fulltext_template_direct 分析。",
                    "caution": "请在深度分析队列中选择 fulltext_template_direct 后重新分析，再导出正式报告。",
                },
                "sections": [],
            }

        include_items = []
        review_items = []
        submm_review_items = []
        limitation_items = []
        exclude_count = 0
        for result, item, payload in self._dedupe_template_direct_rows(payloads):
            for evidence in payload.get("evidences", []):
                if not isinstance(evidence, dict):
                    continue
                recommendation = evidence.get("recommendation")
                row = (result, item, payload, evidence)
                if recommendation == "include":
                    include_items.append(row)
                elif recommendation == "review":
                    if self._is_template_direct_limitation(evidence):
                        limitation_items.append(row)
                    elif self._is_template_direct_submm_candidate(evidence):
                        submm_review_items.append(row)
                    else:
                        review_items.append(row)
                else:
                    exclude_count += 1
                    if include_all:
                        if self._is_template_direct_limitation(evidence):
                            limitation_items.append(row)
                        else:
                            review_items.append(row)
        summary = {
            "session_label": session.display_name if session else session_id,
            "include_count": len(include_items),
            "review_count": len(review_items) + len(submm_review_items),
            "submm_review_count": len(submm_review_items),
            "limitation_count": len(limitation_items),
            "exclude_count": exclude_count,
        }
        conclusion = self._template_direct_conclusion(
            include_items,
            review_items,
            submm_review_items,
            limitation_items,
            exclude_count,
        )
        submm_items = [row for row in include_items if row[3].get("claim_type") == "submm_precision_claim"]
        capability_items = [row for row in include_items if row not in submm_items]
        sections = [
            {"title": "推荐纳入", "items": [], "empty": not include_items},
            {
                "title": "直接亚毫米级佐证",
                "items": self._template_direct_item_models(submm_items, start_index=1),
            },
            {
                "title": "能力认可佐证",
                "items": self._template_direct_item_models(capability_items, start_index=len(submm_items) + 1),
            },
            {
                "title": "直接亚毫米级佐证候选：需人工核对引用编号",
                "items": self._template_direct_item_models(submm_review_items, start_index=len(include_items) + 1),
            },
            {
                "title": "候选复核附录",
                "items": self._template_direct_item_models(
                    review_items,
                    start_index=len(include_items) + len(submm_review_items) + 1,
                ),
            },
            {
                "title": "局限性反馈 / 不宜作为亮点",
                "items": self._template_direct_item_models(
                    limitation_items,
                    start_index=len(include_items) + len(submm_review_items) + len(review_items) + 1,
                ),
            },
            {
                "title": "不纳入证据摘要",
                "items": [],
                "summary_text": f"本次分析中有 {exclude_count} 条证据被模型建议不纳入正式报告。",
            },
        ]
        return {
            "has_results": True,
            "session_label": session.display_name if session else session_id,
            "summary": summary,
            "conclusion": conclusion,
            "sections": sections,
        }

    def _template_direct_conclusion(
        self,
        include_items: list,
        review_items: list,
        submm_review_items: list,
        limitation_items: list,
        exclude_count: int,
    ) -> dict:
        submm_items = [row for row in include_items if row[3].get("claim_type") == "submm_precision_claim"]
        capability_items = [
            row for row in include_items
            if row[3].get("claim_type") in {
                "through_wall_eavesdropping",
                "rfid_loudspeaker_vibration",
                "capability_recognition",
                "method_use",
                "performance_comparison",
            }
        ]
        first_items = [
            row for row in include_items
            if row[3].get("claim_type") in {"first_or_pioneering_claim", "first_or_seminal_claim"}
        ]
        if submm_items or submm_review_items:
            sample_row = (submm_items or submm_review_items)[0]
            sample = self._short_plain_text(sample_row[3].get("evidence_quote") or "", max_chars=180)
            if submm_items:
                submm_line = f"发现少量已核验直接亚毫米级文本证据，例如 “{sample}”。"
            else:
                submm_line = f"发现直接亚毫米级强候选文本证据，例如 “{sample}”，但引用编号或参考文献条目仍需人工核验。"
        else:
            submm_line = "未发现可靠的第三方正文直接亚毫米级评价；如有候选，应人工复核引用编号和原文语义。"
        first_line = (
            "发现明确 first / pioneering 评价。"
            if first_items
            else "目前未发现可靠第三方明确将 first / pioneering 作用到目标论文。"
        )
        capability_line = (
            f"发现 {len(capability_items)} 条能力认可证据，主要涉及 RFID through-wall eavesdropping、speaker/loudspeaker vibration sensing 或明确方法/性能使用。"
            if capability_items
            else "未发现可直接纳入主报告的能力认可证据。"
        )
        return {
            "submm": submm_line,
            "first": first_line,
            "capability": capability_line,
            "caution": f"普通相关工作、成组引用、标题-only、reference-only、局限性反馈不能作为正向亮点；本次另有 {len(submm_review_items)} 条直接亚毫米候选、{len(review_items)} 条候选复核、{len(limitation_items)} 条局限性/不宜作为亮点、{exclude_count} 条不纳入。",
        }

    def _render_template_direct_markdown(self, model: dict) -> str:
        if not model.get("has_results"):
            summary = model["summary"]
            return "\n".join(
                [
                    "# 亮点引用证据报告",
                    "",
                    "## 一、报告摘要",
                    "",
                    f"- 学者会话：{summary['session_label']}",
                    "- 推荐纳入证据数：0",
                    "- 候选复核证据数：0",
                    "- 不纳入证据数：0",
                    "",
                    "尚未运行 fulltext_template_direct 分析。",
                    "",
                    "请在深度分析队列中选择 fulltext_template_direct 后重新分析，再导出正式报告。",
                    "",
                ]
            ).rstrip() + "\n"
        summary = model["summary"]
        conclusion = model["conclusion"]
        lines = [
            "# 亮点引用证据报告",
            "",
            "## 一、报告摘要",
            "",
            f"- 学者会话：{summary['session_label']}",
            f"- 推荐纳入证据数：{summary['include_count']}",
            f"- 候选复核证据数：{summary['review_count']}",
            f"- 局限性/不宜作为亮点证据数：{summary['limitation_count']}",
            f"- 不纳入证据数：{summary['exclude_count']}",
            "",
            "## 结论摘要",
            "",
            f"- 是否发现第三方明确亚毫米级佐证：{conclusion['submm']}",
            f"- 是否发现 first / pioneering 评价：{conclusion['first']}",
            f"- 是否发现能力认可：{conclusion['capability']}",
            f"- 不宜过度解读：{conclusion['caution']}",
            "",
        ]
        section_number = 2
        for section in model["sections"]:
            title = section["title"]
            if title in {"直接亚毫米级佐证", "能力认可佐证"}:
                heading = f"### {title}"
            else:
                heading = f"## {self._chinese_section_number(section_number)}、{title}"
                section_number += 1
            lines.extend([heading, ""])
            if section.get("summary_text"):
                lines.extend([section["summary_text"], ""])
                continue
            items = section.get("items") or []
            if not items:
                lines.extend(["暂无。", ""])
                continue
            for item in items:
                self._append_template_direct_item_markdown(lines, item)
        return "\n".join(lines).rstrip() + "\n"

    def _chinese_section_number(self, value: int) -> str:
        return {1: "一", 2: "二", 3: "三", 4: "四", 5: "五", 6: "六", 7: "七"}.get(value, str(value))

    def _template_direct_item_models(self, items: list, *, start_index: int) -> list:
        return [
            self._template_direct_item_model(row, index)
            for index, row in enumerate(items, start=start_index)
        ]

    def _template_direct_item_model(self, row: tuple, index: int) -> dict:
        _result, item, payload, evidence = row
        marker = evidence.get("evidence_reference_marker") or payload.get("target_reference_marker") or ""
        reference_entry = (
            evidence.get("evidence_reference_entry_raw")
            or f"{marker} 未解析到引用论文原文 References 条目"
            or ""
        )
        claim_type = evidence.get("claim_type") or "ordinary_reference"
        recommendation_label = {
            "include": "推荐纳入",
            "review": "候选复核",
            "exclude": "不纳入",
        }.get(evidence.get("recommendation"), evidence.get("recommendation") or "候选复核")
        quote = self._highlight_template_direct_text(
            evidence.get("evidence_quote") or "",
            marker=marker,
            reference_entry=reference_entry,
            item=item,
        )
        context = self._trim_markdown_context(
            evidence.get("evidence_context") or evidence.get("evidence_quote") or "",
            citation_text=evidence.get("evidence_quote") or "",
        )
        context = self._highlight_template_direct_text(
            context,
            marker=marker,
            reference_entry=reference_entry,
            item=item,
        )
        return {
            "index": index,
            "claim_type": claim_type,
            "claim_label": (
                evidence.get("template_display_label")
                if evidence.get("template_satisfied") and evidence.get("template_display_label")
                else self._template_direct_claim_label(claim_type)
            ),
            "citing_paper_title": item.citing_paper_title if item else "unknown",
            "cited_paper_title": item.cited_paper_title if item else "unknown",
            "recommendation_label": recommendation_label,
            "confidence": evidence.get("confidence") or "unknown",
            "quote_markdown": quote,
            "quote_html": self._highlight_template_direct_html(
                evidence.get("evidence_quote") or "",
                marker=marker,
                reference_entry=reference_entry,
                item=item,
            ),
            "reference_entry_markdown": self._highlight_template_direct_text(
                reference_entry,
                marker=marker,
                reference_entry=reference_entry,
                item=item,
            ),
            "reference_entry_html": self._highlight_template_direct_html(
                reference_entry,
                marker=marker,
                reference_entry=reference_entry,
                item=item,
            ),
            "context_markdown": context,
            "context_html": self._highlight_template_direct_html(
                self._trim_markdown_context(
                    evidence.get("evidence_context") or evidence.get("evidence_quote") or "",
                    citation_text=evidence.get("evidence_quote") or "",
                ),
                marker=marker,
                reference_entry=reference_entry,
                item=item,
            ),
            "copy_ready_zh": self._sanitize_markdown_report_text(evidence.get("copy_ready_zh") or ""),
            "why_this_judgment_zh": self._sanitize_markdown_report_text(evidence.get("why_this_judgment_zh") or ""),
        }

    def _append_template_direct_item_markdown(self, lines: List[str], item: dict) -> None:
        lines.extend(
            [
                f"### {item['index']}. {item['claim_label']}",
                "",
                f"**引用论文：** {item['citing_paper_title']}  ",
                f"**被引用论文：** {item['cited_paper_title']}  ",
                f"**建议：** {item['recommendation_label']}  ",
                f"**置信度：** {item['confidence']}  ",
                "",
                "#### 原文证据",
                "",
                self._blockquote(item["quote_markdown"]),
                "",
                "#### 对应参考文献（引用论文原文 References 中的条目）",
                "",
                self._blockquote(item["reference_entry_markdown"]),
                "",
                "#### 原文上下文",
                "",
                self._blockquote(item["context_markdown"]),
                "",
                "#### 亮点评价",
                "",
                item["copy_ready_zh"],
                "",
                "#### 评价理由",
                "",
                item["why_this_judgment_zh"],
                "",
            ]
        )

    def _is_template_direct_limitation(self, evidence: dict) -> bool:
        if (evidence.get("claim_type") or "") in {"limitation_feedback", "method_use_with_limitation"}:
            return True
        text = " ".join(
            str(evidence.get(key) or "")
            for key in ("evidence_quote", "evidence_context", "why_this_judgment_zh", "copy_ready_zh", "postprocess_reason")
        ).lower()
        return any(
            phrase in text
            for phrase in (
                "less practical",
                "requires pre-installing",
                "pre-installed tag",
                "limited",
                "not practical",
                "impractical",
                "title_or_reference_only_not_include",
                "title-only",
                "reference-only",
                "标题本身",
                "局限性",
            )
        )

    def _is_template_direct_submm_candidate(self, evidence: dict) -> bool:
        return (evidence.get("claim_type") or "") == "submm_precision_claim"

    def _dedupe_template_direct_rows(self, payloads: list) -> list:
        kept = {}
        order = []
        for result, item, payload in payloads:
            for evidence in payload.get("evidences", []):
                if not isinstance(evidence, dict):
                    continue
                key = (
                    item.citing_paper_title if item else "",
                    evidence.get("evidence_reference_marker") or payload.get("target_reference_marker") or "",
                    self._normalize_report_quote(evidence.get("evidence_quote") or ""),
                )
                row = (result, item, payload, evidence)
                if key not in kept:
                    kept[key] = row
                    order.append(key)
                elif self._template_direct_row_rank(row) > self._template_direct_row_rank(kept[key]):
                    kept[key] = row
        grouped = {}
        grouped_order = []
        for key in order:
            result, item, payload, evidence = kept[key]
            group_key = (result.id, item.id if item else None)
            if group_key not in grouped:
                grouped[group_key] = (result, item, dict(payload, evidences=[]))
                grouped_order.append(group_key)
            grouped[group_key][2]["evidences"].append(evidence)
        return [grouped[key] for key in grouped_order]

    def _template_direct_row_rank(self, row: tuple) -> tuple:
        _result, _item, _payload, evidence = row
        recommendation_rank = {"include": 3, "review": 2, "exclude": 1}.get(evidence.get("recommendation"), 0)
        claim_rank = {
            "submm_precision_claim": 90,
            "through_wall_eavesdropping": 80,
            "rfid_loudspeaker_vibration": 70,
            "performance_comparison": 60,
            "custom_template_evidence": 55,
            "method_use": 50,
            "capability_recognition": 40,
            "limitation_feedback": 30,
            "ordinary_reference": 20,
            "false_positive": 0,
        }.get(evidence.get("claim_type"), 0)
        return recommendation_rank, claim_rank

    def _normalize_report_quote(self, value: str) -> str:
        return re.sub(r"\s+", " ", unicodedata.normalize("NFKD", str(value or ""))).strip().casefold()

    def _normalized_template_direct_payload(self, result, item, payload: dict) -> dict:
        if item is None:
            return payload
        cited_publication = self.db.get(Publication, item.cited_publication_id) if getattr(item, "cited_publication_id", None) else None
        candidate_payload = self._load_json(result.candidate_spans_json)
        return postprocess_template_direct_payload(
            payload,
            citing_paper_title=item.citing_paper_title or "",
            cited_paper_title=item.cited_paper_title or "",
            cited_paper_doi=getattr(cited_publication, "doi", None),
            target_reference_marker=payload.get("target_reference_marker") or candidate_payload.get("target_reference_marker", ""),
            target_reference_entry=payload.get("target_reference_entry") or candidate_payload.get("target_reference_entry", ""),
            reference_entries_by_marker=self._reference_entries_for_item(item),
            cited_paper_authors=self._load_json_list(
                getattr(cited_publication, "authors_json", None)
            ),
            cited_paper_year=getattr(cited_publication, "year", None),
            target_reference_resolved=(
                candidate_payload.get("reference_anchor_source")
                == "deterministic_resolver"
            ),
        )

    def _reference_entries_for_item(self, item) -> dict:
        if item is None or not getattr(item, "pdf_asset_id", None):
            return {}
        asset = self.db.get(PdfAsset, item.pdf_asset_id)
        path = Path(asset.extracted_text_path) if asset and asset.extracted_text_path else None
        if not path or not path.exists():
            return {}
        try:
            return reference_entries_by_marker(path.read_text(encoding="utf-8"))
        except OSError:
            return {}

    def _template_direct_claim_label(self, claim_type: str) -> str:
        return {
            "submm_precision_claim": "直接亚毫米精度佐证",
            "capability_recognition": "能力佐证",
            "through_wall_eavesdropping": "穿墙窃听能力佐证",
            "rfid_loudspeaker_vibration": "扬声器振动感知能力佐证",
            "method_use": "方法使用",
            "method_foundation": "方法基础",
            "theoretical_foundation": "理论基础",
            "application_extension": "应用扩展",
            "detailed_comparison": "详细对比",
            "performance_comparison": "性能对比",
            "custom_template_evidence": "自定义模板证据",
            "baseline_or_benchmark": "基线 / Benchmark",
            "positive_evaluation": "正向评价",
            "first_or_seminal_claim": "首次 / 开创性评价",
            "representative_work": "代表性相关工作",
            "background_reference": "背景引用",
            "ordinary_citation": "普通引用",
            "limitation_feedback": "局限性反馈",
            "limitation_or_negative": "局限性反馈",
            "ordinary_reference": "普通相关工作",
            "false_positive": "误报候选",
        }.get(claim_type, claim_type)

    def _highlight_template_direct_text(self, text: str, *, marker: str, reference_entry: str, item) -> str:
        highlights = []
        if marker:
            highlights.append(marker)
        highlights.extend(self._semantic_claim_highlights(text))
        if item is not None and item.cited_paper_title:
            highlights.append(item.cited_paper_title)
            if ":" in item.cited_paper_title:
                highlights.append(item.cited_paper_title.split(":", 1)[0])
        doi_match = re.search(r"\b10\.\d{4,9}/[-._;()/:A-Z0-9]+\b", reference_entry or "", flags=re.I)
        if doi_match:
            highlights.append(doi_match.group(0))
        return render_markdown_highlight(text or "", highlights, marker)

    def _highlight_template_direct_html(
        self,
        text: str,
        *,
        marker: str,
        reference_entry: str,
        item,
        extra_highlights: Optional[List[str]] = None,
    ) -> str:
        source = str(text or "")
        if not source:
            return ""
        terms: List[Tuple[str, str]] = []
        if marker:
            terms.append((marker, "citation-marker"))
        for phrase in self._semantic_claim_highlights(source):
            terms.append((phrase, "claim-phrase"))
        for phrase in extra_highlights or []:
            if str(phrase or "").strip():
                terms.append((str(phrase), "claim-phrase"))
        if item is not None and item.cited_paper_title:
            terms.append((item.cited_paper_title, "target-reference"))
            if ":" in item.cited_paper_title:
                terms.append((item.cited_paper_title.split(":", 1)[0], "target-reference"))
        doi_match = re.search(r"\b10\.\d{4,9}/[-._;()/:A-Z0-9]+\b", reference_entry or "", flags=re.I)
        if doi_match:
            terms.append((doi_match.group(0), "target-reference"))

        spans: List[Tuple[int, int, str]] = []
        seen_terms = set()
        for term, class_name in sorted(terms, key=lambda pair: len(pair[0]), reverse=True):
            value = str(term or "").strip()
            if not value:
                continue
            key = (value.casefold(), class_name)
            if key in seen_terms:
                continue
            seen_terms.add(key)
            for match in re.finditer(re.escape(value), source, flags=re.I):
                start, end = match.span()
                if any(not (end <= existing_start or start >= existing_end) for existing_start, existing_end, _ in spans):
                    continue
                spans.append((start, end, class_name))
        if not spans:
            return html.escape(source)
        parts: List[str] = []
        cursor = 0
        for start, end, class_name in sorted(spans):
            parts.append(html.escape(source[cursor:start]))
            parts.append(f'<mark class="{class_name}">{html.escape(source[start:end])}</mark>')
            cursor = end
        parts.append(html.escape(source[cursor:]))
        return "".join(parts)

    def _semantic_claim_highlights(self, text: str) -> List[str]:
        source = str(text or "")
        patterns = [
            r"sub\s*-\s*millimeter\s*-\s*level\s+vibrations?",
            r"sub\s*-\s*mm\s+level\s+vibrations?",
            r"submillimeter",
            r"through\s*-\s*the\s*-\s*wall\s+sound\s+eavesdropping",
            r"thru\s*-\s*the\s*-\s*wall\s+sound\s+eavesdropping",
            r"through\s*-\s*wall\s+eavesdropping",
            r"thru\s*-\s*wall\s+eavesdropping",
            r"RFID\s+tags?",
            r"loudspeaker\s+vibrations?",
            r"speaker\s+vibrations?",
            r"vibrations?\s+from\s+loudspeakers?",
            r"reconstructs?\s+audio",
        ]
        phrases: List[str] = []
        seen = set()
        for pattern in patterns:
            for match in re.finditer(pattern, source, flags=re.I):
                phrase = match.group(0)
                key = phrase.casefold()
                if key not in seen:
                    phrases.append(phrase)
                    seen.add(key)
        return phrases

    def _inline_markdown_to_html(self, value: str) -> str:
        escaped = html.escape(str(value or ""))
        return re.sub(r"\*\*(.+?)\*\*", r"<mark>\1</mark>", escaped)

    def export_pptx(self, session_id: int) -> Path:
        path = self._session_export_dir(session_id) / "report.pptx"
        diagnostics_path = self._session_export_dir(session_id) / "report_pptx_diagnostics.json"
        cards = [
            card for card in self.list_cards(session_id)
            if card.include_in_report and card.review_status not in {"false_positive", "rejected"}
        ]
        warnings: List[str] = []
        for card in cards:
            for image_path in self._candidate_image_paths(card):
                path_obj = Path(image_path)
                if not path_obj.exists() or not path_obj.is_file() or path_obj.stat().st_size == 0:
                    warnings.append(f"missing_image:{path_obj.name}")
        temp_fd, temp_name = tempfile.mkstemp(suffix=".pptx", dir=str(self._session_export_dir(session_id)))
        os.close(temp_fd)
        Path(temp_name).unlink(missing_ok=True)
        temp_path = Path(temp_name)
        try:
            self._write_minimal_pptx(temp_path, session_id, cards)
            validation_error = self._validate_pptx(temp_path)
            if validation_error is not None:
                diagnostics_path.write_text(
                    json.dumps(
                        {
                            "export_status": "failed",
                            "slide_count": len(cards) + 1,
                            "warning_count": len(warnings),
                            "warnings": warnings,
                            "validation_status": "failed",
                            "validation_error": validation_error,
                        },
                        ensure_ascii=False,
                        indent=2,
                    ),
                    encoding="utf-8",
                )
                temp_path.unlink(missing_ok=True)
                raise PptxExportError(f"PPTX 导出失败：{validation_error}")
            shutil.move(str(temp_path), str(path))
            diagnostics_path.write_text(
                json.dumps(
                    {
                        "export_status": "succeeded",
                        "slide_count": len(cards) + 2,
                        "warning_count": len(warnings),
                        "warnings": warnings,
                        "validation_status": "ok",
                        "validation_error": "",
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
                encoding="utf-8",
            )
        finally:
            temp_path.unlink(missing_ok=True)
        return path

    def write_cards_csv(self, session_id: int) -> Path:
        path = self._session_export_dir(session_id) / "highlight_cards.csv"
        path.write_text(self.export_cards_csv(session_id), encoding="utf-8")
        return path

    def write_cards_markdown(self, session_id: int) -> Path:
        path = self._session_export_dir(session_id) / "highlight_cards.md"
        path.write_text(self.export_cards_markdown(session_id), encoding="utf-8")
        return path

    def _should_export_markdown_row(self, row: dict, *, include_all: bool) -> bool:
        card = row["card"]
        meta = row.get("narrative_meta") or {}
        if include_all:
            return True
        if card.strong_evidence_id:
            evidence = self.db.get(StrongEvidence, card.strong_evidence_id)
            if evidence is None:
                return False
            if (evidence.review_status or "") in {"false_positive", "rejected"}:
                return False
            if not (evidence.citation_text or "").strip():
                return False
        if not getattr(card, "include_in_report", False):
            return False
        if (card.review_status or "") in {"false_positive", "rejected"}:
            return False
        if meta.get("anchor_validation_status") == "invalid":
            return False
        if meta.get("anchor_validation_reason") == "cited_other_reference_marker":
            return False
        return True

    def _report_card_view_model(self, index: int, row: dict) -> ReportCardViewModel:
        card = row["card"]
        context_preview = row.get("context_preview") or {}
        meta = row.get("narrative_meta") or {}
        evidence = self.db.get(StrongEvidence, card.strong_evidence_id) if card.strong_evidence_id else None
        item = self.db.get(DeepAnalysisQueueItem, evidence.queue_item_id) if evidence and evidence.queue_item_id else None
        citation_text = (
            meta.get("evidence_quote")
            or card.evidence_quote
            or (evidence.citation_text if evidence else "")
            or ""
        )
        key_phrases = self._markdown_key_phrases(meta=meta, context_preview=context_preview, evidence=evidence)
        target_marker = (
            meta.get("target_reference_marker")
            or context_preview.get("target_reference_marker")
            or ""
        )
        display_context = self._markdown_display_context(
            context_preview=context_preview,
            meta=meta,
            citation_text=citation_text,
        )
        why = self._markdown_why_this_judgment(
            card=card,
            meta=meta,
            key_phrases=key_phrases,
            target_marker=target_marker,
            section_heading=context_preview.get("section_heading") or "",
        )
        statement = self._markdown_copy_statement(card=card, meta=meta, why=why)
        risk_note = (
            meta.get("risk_note")
            or ("如该段为成组引用，需要人工确认归因范围。" if card.review_status in {"unreviewed", "needs_discussion"} else "")
            or "未发现额外风险提示；仍建议人工复核原文上下文后使用。"
        )
        return ReportCardViewModel(
            card_index=index,
            card_type_label_zh=self._card_type_label(card.card_type),
            title_zh=card.title or self._card_type_label(card.card_type),
            citing_paper_title=card.source_citing_paper_title or "",
            citing_venue=card.venue or (item.venue if item else "") or "unknown",
            citing_year=str(item.year or "") if item and item.year else "unknown",
            cited_paper_title=card.source_cited_paper_title or "",
            evidence_type_label_zh=self._card_type_label(card.aspect or card.card_type),
            evidence_strength=card.evidence_strength or "unknown",
            stance=card.stance or "unknown",
            notable_author_name=card.notable_author_name or "",
            honor_category="" if (card.fellow_status or "unknown") == "unknown" else (card.fellow_status or ""),
            citation_text=citation_text,
            display_context=display_context,
            highlighted_citation_markdown=render_markdown_highlight(
                citation_text,
                key_phrases,
                target_marker,
            ),
            highlighted_context_markdown=render_markdown_highlight(
                display_context,
                key_phrases,
                target_marker,
            ),
            key_phrases=key_phrases,
            target_reference_marker=target_marker,
            why_this_judgment=self._sanitize_markdown_report_text(why),
            limitation_zh=self._sanitize_markdown_report_text(
                meta.get("limitation_zh") or "请严格按照原文上下文表述，不要扩展为原文没有支持的评价。"
            ),
            copy_ready_statement=self._sanitize_markdown_report_text(statement),
            report_recommendation=meta.get("report_recommendation") or "",
            risk_note=self._sanitize_markdown_report_text(risk_note),
            confidence_level=meta.get("confidence_level") or meta.get("confidence") or "",
            anchor_validation_status=meta.get("anchor_validation_status") or "unknown",
            anchor_validation_reason=meta.get("anchor_validation_reason") or "",
            include_in_report=bool(card.include_in_report),
            review_status=card.review_status or "unreviewed",
            matched_template_names=meta.get("matched_template_names") or [],
            template_match_reason=meta.get("template_match_reason") or "",
            template_satisfied=meta.get("template_satisfied"),
            template_failure_reason=meta.get("template_failure_reason") or "",
        )

    def _append_markdown_card_section(
        self,
        lines: List[str],
        title: str,
        cards: List[ReportCardViewModel],
    ) -> None:
        lines.extend([title, ""])
        if not cards:
            lines.extend(["暂无。", ""])
            return
        seen = set()
        for card in cards:
            key = (card.title_zh, card.citation_text, card.card_type_label_zh)
            if key in seen:
                continue
            seen.add(key)
            lines.extend(self._render_markdown_card(card))

    def _render_markdown_card(self, card: ReportCardViewModel) -> List[str]:
        review_advice = "需要复核" if card.review_status in {"unreviewed", "needs_discussion"} or card.risk_note else "可纳入报告"
        heading = f"### {card.card_index}. {card.card_type_label_zh}：{card.title_zh}"
        context_note = ""
        if card.display_context.strip() == card.citation_text.strip():
            context_note = "\n\n> 未找到更长上下文，仅展示核心证据句。"
        notable_lines = []
        if card.notable_author_name:
            notable = card.notable_author_name
            if card.honor_category:
                notable = f"{notable}（{card.honor_category}）"
            notable_lines.append(f"**重要作者：** {notable}  ")
        return [
            heading,
            "",
            f"**引用论文：** {card.citing_paper_title}  ",
            f"**发表位置：** {card.citing_venue}, {card.citing_year}  ",
            f"**被引用论文：** {card.cited_paper_title}  ",
            f"**证据类型：** {card.evidence_type_label_zh} / {card.card_type_label_zh}  ",
            f"**证据强度：** {card.evidence_strength}  ",
            f"**置信度：** {card.confidence_level or 'unknown'}  ",
            f"**人工复核建议：** {review_advice}  ",
            f"**报告推荐：** {card.report_recommendation or '候选复核'}  ",
            f"**命中模板：** {', '.join(card.matched_template_names) if card.matched_template_names else '无'}  ",
            f"**模板是否满足：** {card.template_satisfied if card.template_satisfied is not None else 'unknown'}  ",
            *notable_lines,
            "",
            "### 原文证据",
            "",
            self._blockquote(card.highlighted_citation_markdown),
            "",
            "### 原文上下文",
            "",
            self._blockquote(card.highlighted_context_markdown) + context_note,
            "",
            "### 亮点评价",
            "",
            card.copy_ready_statement,
            "",
            "### 评价理由",
            "",
            card.why_this_judgment,
            "",
            "### 不应如何过度解读",
            "",
            card.limitation_zh,
            "",
            "### 模板判断",
            "",
            self._template_markdown_summary(card),
            "",
            "### 风险提示",
            "",
            card.risk_note,
            "",
            "<!-- "
            f"Citing paper: {card.citing_paper_title}; "
            f"Evidence quote: {card.citation_text}; "
            f"Evidence reason: {card.why_this_judgment}; "
            f"card_type: {card.card_type_label_zh}; "
            f"anchor_validation_status: {card.anchor_validation_status}; "
            f"anchor_validation_reason: {card.anchor_validation_reason}"
            f"; matched_templates: {', '.join(card.matched_template_names)}"
            f"; template_match_reason: {card.template_match_reason}"
            f"; template_failure_reason: {card.template_failure_reason}"
            " -->",
            "",
        ]

    def _template_markdown_summary(self, card: ReportCardViewModel) -> str:
        if card.template_satisfied:
            return (
                f"命中模板：{', '.join(card.matched_template_names) or '未记录名称'}。\n\n"
                f"模板命中原因：{card.template_match_reason or '未记录原因'}。"
            )
        if card.template_failure_reason:
            return f"未满足当前模板：{card.template_failure_reason}"
        return "未记录模板命中信息。"

    def _markdown_key_phrases(self, *, meta: dict, context_preview: dict, evidence: Optional[StrongEvidence]) -> List[str]:
        phrases: List[str] = []
        for collection in [
            meta.get("key_phrases") or [],
            meta.get("technical_terms_used") or [],
            context_preview.get("highlight_terms") or [],
            self._load_json_list(evidence.highlight_keywords_json) if evidence else [],
        ]:
            for phrase in collection:
                value = str(phrase or "").strip()
                if value and value.casefold() not in {item.casefold() for item in phrases}:
                    phrases.append(value)
        return phrases[:10]

    def _markdown_display_context(self, *, context_preview: dict, meta: dict, citation_text: str) -> str:
        context = (
            context_preview.get("display_context")
            or meta.get("evidence_context")
            or context_preview.get("body_context_full")
            or citation_text
            or ""
        )
        return self._trim_markdown_context(str(context), citation_text=citation_text)

    def _trim_markdown_context(self, context: str, *, citation_text: str, max_chars: int = 1500) -> str:
        text = re.sub(r"\s+", " ", context or "").strip()
        if len(text) <= max_chars:
            return text
        needle = (citation_text or "").strip()
        index = text.find(needle) if needle else -1
        if index < 0:
            return text[:max_chars].rsplit(" ", 1)[0] + " ..."
        half = max_chars // 2
        start = max(0, index - half)
        end = min(len(text), index + len(needle) + half)
        excerpt = text[start:end].strip()
        if start > 0:
            excerpt = "... " + excerpt
        if end < len(text):
            excerpt = excerpt + " ..."
        return excerpt

    def _markdown_why_this_judgment(
        self,
        *,
        card: HighlightCard,
        meta: dict,
        key_phrases: List[str],
        target_marker: str,
        section_heading: str,
    ) -> str:
        existing = str(meta.get("why_this_judgment") or "").strip()
        if existing and len(existing) > 12:
            return existing
        label = self._card_type_label(card.card_type)
        marker_text = f"正文包含目标引用编号 {target_marker}" if target_marker else "正文证据需要结合目标论文标题或作者信息复核"
        phrase_text = "、".join(key_phrases[:5]) if key_phrases else "原文证据中的关键表述"
        section_text = f"该段位于 {section_heading}，" if section_heading else "该段来自引用论文正文，"
        grouped = "成组引用" in (card.narrative_zh or "") or card.review_status in {"unreviewed", "needs_discussion"}
        risk = "；但该证据可能涉及成组引用，需要人工确认归因范围" if grouped else ""
        return (
            f"{section_text}{marker_text}，并围绕 {phrase_text} 展开表述。"
            f"这些原文线索支持将其判断为“{label}”，而不是只依据字段标签自动归类{risk}。"
        )

    def _markdown_copy_statement(self, *, card: HighlightCard, meta: dict, why: str) -> str:
        statement = str(meta.get("copy_ready_statement") or card.narrative_zh or card.body_markdown or "").strip()
        if statement:
            return statement
        label = self._card_type_label(card.card_type)
        return (
            f"引用论文《{card.source_citing_paper_title}》在正文中引用了目标论文《{card.source_cited_paper_title}》，"
            f"可作为“{label}”类引用证据候选。{why}"
        )

    def _sanitize_markdown_report_text(self, value: str) -> str:
        return str(value or "").replace("高度评价", "直接正向赞扬")

    def _short_plain_text(self, value: str, *, max_chars: int = 180) -> str:
        text = re.sub(r"\s+", " ", str(value or "")).strip()
        if len(text) <= max_chars:
            return text
        return text[:max_chars].rsplit(" ", 1)[0] + " ..."

    def _blockquote(self, text: str) -> str:
        value = str(text or "").strip()
        if not value:
            return "> （无可导出的原文内容）"
        return "\n".join(f"> {line}" if line else ">" for line in value.splitlines())

    def _excluded_false_positive_summary(self, excluded_rows: List[dict]) -> str:
        if not excluded_rows:
            return "默认报告未排除误报卡片。"
        counts = {}
        for row in excluded_rows:
            reason = (row.get("narrative_meta") or {}).get("anchor_validation_reason") or "false_positive"
            counts[reason] = counts.get(reason, 0) + 1
        return "；".join(f"{reason}: {count}" for reason, count in sorted(counts.items()))

    def _session_export_dir(self, session_id: int) -> Path:
        path = Path(settings.export_dir) / f"scholar_session_{session_id}"
        path.mkdir(parents=True, exist_ok=True)
        return path

    def _find_notable_author_for_item(
        self,
        queue_item_id: int,
        citing_authors_json: Optional[str],
    ) -> Optional[NotableAuthor]:
        annotation = (
            self.db.query(CitationAuthorAnnotation)
            .filter(
                CitationAuthorAnnotation.queue_item_id == queue_item_id,
                CitationAuthorAnnotation.is_important.is_(True),
                CitationAuthorAnnotation.match_status == "matched",
            )
            .order_by(CitationAuthorAnnotation.id.asc())
            .first()
        )
        if annotation is not None:
            notable = self.db.get(NotableAuthor, annotation.notable_author_id)
            if notable is not None:
                return notable
        return self._find_notable_author(citing_authors_json)

    def _fallback_workspace_rows(
        self,
        session_id: int,
        covered_queue_item_ids: set,
        view: str = "all",
    ) -> List[dict]:
        from app.models import DeepAnalysisQueueItem, FulltextAnalysisResult

        rows = []
        items = (
            self.db.query(DeepAnalysisQueueItem)
            .filter_by(scholar_session_id=session_id, queue_status="analyzed")
            .order_by(DeepAnalysisQueueItem.id.asc())
            .all()
        )
        for item in items:
            if item.id in covered_queue_item_ids:
                continue
            latest_result = (
                self.db.query(FulltextAnalysisResult)
                .filter_by(queue_item_id=item.id, status="succeeded")
                .order_by(FulltextAnalysisResult.id.desc())
                .first()
            )
            if latest_result is None:
                continue
            diagnostics = self._load_json(latest_result.candidate_spans_json)
            parsed = self._load_json(latest_result.parsed_result_json)
            findings = parsed.get("findings", []) if isinstance(parsed, dict) else []
            card = self._synthetic_fallback_card(item, diagnostics, findings)
            validation = validate_citation_target_anchor(
                citation_text=card.evidence_quote or "",
                target_reference_marker=diagnostics.get("target_reference_marker"),
                cited_paper_title=item.cited_paper_title,
                cited_authors_json=item.cited_authors_json,
            )
            if validation.anchor_validation_status != "unknown" and not validation.is_valid:
                if not self._is_false_positive_view(view):
                    continue
                card.review_status = "false_positive"
                card.evidence_strength = "none"
                card.score = 0
                card.title = self._false_positive_title(card.title)
            rows.append(
                {
                    "card": card,
                    "card_kind": "false_positive" if card.review_status == "false_positive" else "ordinary",
                    "grouped_warning": False,
                    "context_preview": self._context_preview_from_result(
                        item=item,
                        result=latest_result,
                        citation_text=card.evidence_quote,
                    ),
                    "narrative_meta": {
                        "risk_note": "这是普通引用/背景引用素材卡片，默认不纳入报告，建议人工复核。",
                        "technical_terms_used": [],
                        "evidence_basis": "",
                        "target_reference_marker": validation.target_reference_marker,
                        "citation_text_contains_target_marker": validation.citation_text_contains_target_marker,
                        "citation_text_contains_other_marker": validation.citation_text_contains_other_marker,
                        "anchor_validation_status": validation.anchor_validation_status,
                        "anchor_validation_reason": validation.anchor_validation_reason,
                    },
                }
            )
        return rows

    def _synthetic_fallback_card(self, item, diagnostics, findings):
        quote = ""
        card_type = "citation_only"
        narrative = "该论文引用了目标论文，但当前未识别出可直接纳入亮点评价的强证据，建议人工复核。"
        if findings:
            first = findings[0]
            quote = str(first.get("citation_text") or "")
            if first.get("evidence_type") == "background":
                if "[36]" in quote or "[60]" in quote or "related work" in quote.lower():
                    card_type = "representative_work"
                else:
                    card_type = "background_reference"
                narrative = (
                    "该论文在正文中将目标论文作为相关工作进行引用，说明目标论文进入了该方向后续研究的技术脉络，但不等同于直接正向评价。"
                )
            else:
                card_type = "weak_mention"
                narrative = (
                    "该论文在正文中引用了目标论文，但当前系统未识别出明确评价、方法采用或详细对比证据，建议人工复核。"
                )
        elif diagnostics.get("target_reference_context_count", 0) > 0:
            card_type = "ordinary_citation"
            preview = diagnostics.get("target_contexts_preview", [])
            if preview:
                quote = str(preview[0].get("context_text_preview") or "")
            narrative = (
                "该论文在正文中引用了目标论文，但当前系统未识别出明确评价、方法采用或详细对比证据，建议人工复核。"
            )
        return SimpleNamespace(
            id=None,
            strong_evidence_id=None,
            source_evidence_id=None,
            scholar_session_id=item.scholar_session_id,
            card_type=card_type,
            title=f"{card_type}: {item.cited_paper_title}",
            subtitle=f"来源论文：{item.citing_paper_title}",
            narrative_zh=narrative,
            body_markdown=narrative,
            evidence_quote=quote,
            highlighted_quote_html=quote,
            source_citing_paper_title=item.citing_paper_title,
            source_cited_paper_title=item.cited_paper_title,
            aspect="background" if card_type in {"background_reference", "citation_only"} else "weak_mention",
            stance="neutral",
            evidence_strength="weak",
            score=0.2,
            user_note="",
            include_in_report=False,
            fellow_status="unknown",
            notable_author_name="",
            notable_author_affiliation="",
            notable_author_role="",
            review_status="unreviewed",
        )

    def _load_json(self, value: Optional[str]):
        if not value:
            return {}
        try:
            return json.loads(value)
        except json.JSONDecodeError:
            return {}

    def _context_preview_for_card(self, card: HighlightCard):
        if not card.strong_evidence_id:
            item = (
                self.db.query(DeepAnalysisQueueItem)
                .filter(DeepAnalysisQueueItem.scholar_session_id == card.scholar_session_id)
                .filter(DeepAnalysisQueueItem.citing_paper_title == card.source_citing_paper_title)
                .filter(DeepAnalysisQueueItem.cited_paper_title == card.source_cited_paper_title)
                .order_by(DeepAnalysisQueueItem.id.asc())
                .first()
            )
            result = None
            if item is not None:
                result = (
                    self.db.query(FulltextAnalysisResult)
                    .filter(FulltextAnalysisResult.queue_item_id == item.id)
                    .order_by(FulltextAnalysisResult.id.desc())
                    .first()
                )
            if item is not None and result is not None:
                return self._context_preview_from_result(item=item, result=result, citation_text=card.evidence_quote or "")
            preview = build_context_preview(extracted_text_path=None, citation_text=card.evidence_quote or "")
            preview["display_context"] = card.evidence_quote or ""
            preview["body_context_full"] = card.evidence_quote or ""
            return preview
        evidence = self.db.get(StrongEvidence, card.strong_evidence_id)
        return self._context_preview_for_evidence(evidence)

    def _context_preview_for_evidence(self, evidence: Optional[StrongEvidence]):
        if evidence is None:
            return build_context_preview(extracted_text_path=None, citation_text="")
        result = self.db.get(FulltextAnalysisResult, evidence.fulltext_result_id)
        if result is None or result.queue_item_id is None:
            return build_context_preview(extracted_text_path=None, citation_text="")
        item = self.db.get(DeepAnalysisQueueItem, result.queue_item_id)
        if item is None or not item.pdf_asset_id:
            return build_context_preview(extracted_text_path=None, citation_text="")
        asset = self.db.get(PdfAsset, item.pdf_asset_id)
        return build_context_preview(
            extracted_text_path=asset.extracted_text_path if asset else None,
            citation_text=evidence.citation_text or "",
            diagnostics=result.candidate_spans_json,
            target_reference_marker=self._load_json(result.candidate_spans_json).get("target_reference_marker"),
            highlight_terms=self._load_json_list(evidence.highlight_keywords_json),
        )

    def _context_preview_from_result(self, *, item, result, citation_text: str):
        if not item.pdf_asset_id:
            return {
                "body_context_full": citation_text,
                "display_context": citation_text,
                "section_heading": "",
                "context_start": None,
                "context_end": None,
            }
        asset = self.db.get(PdfAsset, item.pdf_asset_id)
        return build_context_preview(
            extracted_text_path=asset.extracted_text_path if asset else None,
            citation_text=citation_text,
            diagnostics=result.candidate_spans_json,
            target_reference_marker=self._load_json(result.candidate_spans_json).get("target_reference_marker"),
        )

    def _find_notable_author(self, citing_authors_json: Optional[str]) -> Optional[NotableAuthor]:
        authors = self._load_json_list(citing_authors_json)
        if not authors:
            return None
        for author_name in authors:
            notable = (
                self.db.query(NotableAuthor)
                .filter(NotableAuthor.name == str(author_name))
                .first()
            )
            if notable is not None:
                return notable
        return None

    def _load_json_list(self, value: Optional[str]) -> List[str]:
        if not value:
            return []
        try:
            parsed = json.loads(value)
        except json.JSONDecodeError:
            return []
        return [str(item) for item in parsed] if isinstance(parsed, list) else []

    def _write_minimal_pptx(self, path: Path, session_id: int, cards: List[HighlightCard]) -> None:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        self._add_title_slide(presentation, session_id, cards)
        self._add_summary_slide(presentation, cards)
        for card in cards:
            self._add_card_slide(presentation, card)
        presentation.save(path)

    def _build_slide_texts(self, session_id: int, cards: List[HighlightCard]) -> List[dict]:
        important_count = sum(1 for card in cards if card.review_status == "important")
        notable_count = sum(
            1 for card in cards if (card.fellow_status or "unknown") != "unknown"
        )
        slides = [
            {
                "title": "学术影响力亮点评价汇报",
                "body": _build_slide_body(
                    [
                        f"学者分析会话: {session_id}",
                        f"总证据数: {len(cards)}",
                        f"高价值证据数: {important_count}",
                        f"重要引用团队数: {notable_count}",
                    ]
                ),
            }
        ]
        for card in cards:
            slides.append(
                {
                    "title": card.title,
                    "body": _build_slide_body(
                        [
                            f"引用论文: {card.source_citing_paper_title}",
                            f"被引用论文: {card.source_cited_paper_title}",
                            f"重要作者/团队: {card.notable_author_name or 'unknown'}",
                            f"Fellow 状态: {card.fellow_status or 'unknown'}",
                            f"中文亮点评价: {card.narrative_zh or card.body_markdown}",
                            f"英文原文 citation_text: {self._sanitize_ppt_text(card.evidence_quote)}",
                            f"evidence_type: {card.aspect or ''}",
                            f"stance: {card.stance or ''}",
                            f"score: {card.score or ''}",
                        ]
                    ),
                }
            )
        return slides

    def _candidate_image_paths(self, card: HighlightCard) -> List[str]:
        return []

    def _safe_picture_path(self, image_path: str) -> Optional[Path]:
        path = Path(image_path)
        if not path.exists() or not path.is_file() or path.stat().st_size == 0:
            return None
        if path.suffix.lower() not in {".png", ".jpg", ".jpeg"}:
            return None
        try:
            with Image.open(path) as image:
                image.verify()
        except Exception:
            return None
        return path

    def _sanitize_ppt_text(self, value: Optional[str], limit: int = 1800) -> str:
        text = str(value or "")
        text = "".join(
            ch for ch in text
            if ch in "\n\r\t" or ord(ch) >= 32
        )
        return text[:limit]

    def _validate_pptx(self, path: Path) -> Optional[str]:
        try:
            with zipfile.ZipFile(path) as archive:
                bad = archive.testzip()
                if bad is not None:
                    return f"Bad zip entry: {bad}"
            Presentation(path)
        except Exception as exc:
            return str(exc)
        return None

    def _best_matched_template_type(self, evidence_id: int) -> Optional[str]:
        matches = TemplateService(self.db).list_matches_for_evidence(evidence_id)
        if not matches:
            return None
        best_match = max(matches, key=lambda match: match.match_score)
        template = self.db.get(AnalysisTemplate, best_match.template_id)
        return template.template_type if template is not None else None

    def _template_card_values(self, evidence: StrongEvidence) -> dict:
        template_ids = self._load_json_list(evidence.matched_template_ids_json)
        template_names = []
        for value in template_ids:
            try:
                template = self.db.get(AnalysisTemplate, int(value))
            except (TypeError, ValueError):
                template = None
            if template is not None:
                template_names.append(template.description or template.name)
        return {
            "matched_template_ids_json": json.dumps(template_ids, ensure_ascii=False),
            "matched_template_names": json.dumps(template_names, ensure_ascii=False),
            "template_match_reason": evidence.template_match_reason or "",
            "template_satisfied": evidence.template_satisfied,
            "template_failure_reason": evidence.template_failure_reason or "",
        }

    def _dedupe_evidence_rows(self, rows):
        groups = {}
        for evidence, item in rows:
            card_type = self._best_matched_template_type(evidence.id)
            if not card_type or card_type == "custom":
                card_type = card_type_for_evidence(evidence)
            group_key = (item.id, card_type)
            equivalent_rows = groups.setdefault(group_key, [])
            equivalent_index = next(
                (
                    index
                    for index, (current, _current_item) in enumerate(equivalent_rows)
                    if self._quotes_equivalent(
                        current.citation_text or "",
                        evidence.citation_text or "",
                    )
                ),
                None,
            )
            if equivalent_index is None:
                equivalent_rows.append((evidence, item))
                continue
            current, _current_item = equivalent_rows[equivalent_index]
            if self._prefer_evidence(evidence, current):
                equivalent_rows[equivalent_index] = (evidence, item)
        return [
            row
            for equivalent_rows in groups.values()
            for row in equivalent_rows
        ]

    def _business_key_for_evidence(self, queue_item_id: int, card_type: str, citation_text: str) -> tuple:
        return queue_item_id, card_type, self._canonical_quote_key(citation_text)

    def _canonical_quote_key(self, citation_text: str) -> str:
        text = unicodedata.normalize("NFKD", citation_text or "")
        text = "".join(ch for ch in text if not unicodedata.combining(ch)).lower()
        text = re.sub(r"[^a-z0-9\s]+", " ", text)
        text = " ".join(text.split())
        if len(text) > 220:
            text = text[:220]
        return text

    def _prefer_evidence(self, candidate: StrongEvidence, existing: StrongEvidence) -> bool:
        candidate_key = (
            1 if candidate.review_status == "important" else 0,
            1 if candidate.review_status == "accepted" else 0,
            candidate.score or 0,
            len(candidate.citation_text or ""),
            candidate.id,
        )
        existing_key = (
            1 if existing.review_status == "important" else 0,
            1 if existing.review_status == "accepted" else 0,
            existing.score or 0,
            len(existing.citation_text or ""),
            existing.id,
        )
        return candidate_key > existing_key

    def _select_keeper_card(
        self,
        existing_cards: List[HighlightCard],
        evidence: StrongEvidence,
        card_type: str,
    ) -> Optional[HighlightCard]:
        matches = []
        for card in existing_cards:
            if card.scholar_session_id != evidence.scholar_session_id or card.card_type != card_type:
                continue
            card_evidence = self.db.get(StrongEvidence, card.strong_evidence_id) if card.strong_evidence_id else None
            if card_evidence is None:
                continue
            if card_evidence.queue_item_id != evidence.queue_item_id:
                continue
            if self._quotes_equivalent(card.evidence_quote or "", evidence.citation_text or ""):
                matches.append(card)
        if not matches:
            exact = self.repository.find_by_evidence_id(evidence.id)
            return exact
        matches.sort(key=lambda card: (1 if card.is_user_edited else 0, -card.id), reverse=True)
        return matches[0]

    def _quotes_equivalent(self, left: str, right: str) -> bool:
        a = self._canonical_quote_key(left)
        b = self._canonical_quote_key(right)
        if not a or not b:
            return False
        if a == b or a in b or b in a:
            return True
        shorter, longer = (a, b) if len(a) <= len(b) else (b, a)
        common = sum(1 for token in shorter.split() if token in longer.split())
        return common >= max(5, len(shorter.split()) - 2)

    def _deduplicate_existing_cards(self, session_id: int) -> None:
        cards = self.repository.list_cards(session_id)
        duplicates = []
        keepers: List[HighlightCard] = []
        for card in cards:
            evidence = self.db.get(StrongEvidence, card.strong_evidence_id) if card.strong_evidence_id else None
            if evidence is None:
                keepers.append(card)
                continue
            equivalent_keeper = None
            for keeper in keepers:
                keeper_evidence = self.db.get(StrongEvidence, keeper.strong_evidence_id) if keeper.strong_evidence_id else None
                if keeper_evidence is None:
                    continue
                if keeper.card_type != card.card_type:
                    continue
                if keeper_evidence.queue_item_id != evidence.queue_item_id:
                    continue
                if self._quotes_equivalent(keeper.evidence_quote or "", card.evidence_quote or ""):
                    equivalent_keeper = keeper
                    break
            if equivalent_keeper is None:
                keepers.append(card)
                continue
            preferred = self._prefer_card_keeper(card, equivalent_keeper)
            if preferred is card:
                duplicates.append(equivalent_keeper)
                keepers = [card if existing.id == equivalent_keeper.id else existing for existing in keepers]
            else:
                duplicates.append(card)
        for duplicate in duplicates:
            self.db.delete(duplicate)

    def _mark_mismatched_cards_false_positive(self, session_id: int) -> None:
        changed = False
        for card in self.repository.list_cards(session_id):
            evidence = self.db.get(StrongEvidence, card.strong_evidence_id) if card.strong_evidence_id else None
            if evidence is None or evidence.queue_item_id is None:
                continue
            item = self.db.get(DeepAnalysisQueueItem, evidence.queue_item_id)
            if item is None:
                continue
            result = self.db.get(FulltextAnalysisResult, evidence.fulltext_result_id)
            diagnostics = self._load_json(result.candidate_spans_json if result else None)
            validation = validate_citation_target_anchor(
                citation_text=evidence.citation_text or card.evidence_quote or "",
                target_reference_marker=diagnostics.get("target_reference_marker"),
                cited_paper_title=item.cited_paper_title,
                cited_authors_json=item.cited_authors_json,
            )
            if validation.is_valid or validation.anchor_validation_status == "unknown":
                continue
            self._mark_evidence_and_card_false_positive(evidence, card, validation)
            changed = True
        if changed:
            self.db.commit()

    def _mark_evidence_and_card_false_positive(self, evidence, card, validation) -> None:
        evidence.review_status = "false_positive"
        evidence.anchor_status = validation.anchor_validation_status
        evidence.evidence_strength = "none"
        evidence.score = 0
        note = f"anchor_validation: {validation.anchor_validation_reason}"
        evidence.user_note = self._append_note(evidence.user_note, note)
        if card is None:
            return
        card.review_status = "false_positive"
        card.include_in_report = False
        card.evidence_strength = "none"
        card.score = 0
        card.title = self._false_positive_title(card.title)
        card.user_note = self._append_note(card.user_note, note)

    def _append_note(self, current: Optional[str], note: str) -> str:
        current = (current or "").strip()
        if note in current:
            return current
        return f"{current} {note}".strip()

    def _false_positive_title(self, title: str) -> str:
        title = title or ""
        return title if title.startswith("误报候选：") else f"误报候选：{title}"

    def _is_false_positive_view(self, view: str) -> bool:
        return (view or "all") in {"false_positive", "debug", "invalid_anchor"}

    def _is_false_positive_card(self, card: HighlightCard) -> bool:
        return (card.review_status or "") == "false_positive" or self._is_invalid_anchor_card(card)

    def _is_invalid_anchor_card(self, card: HighlightCard) -> bool:
        evidence = self.db.get(StrongEvidence, card.strong_evidence_id) if card.strong_evidence_id else None
        if evidence is None or evidence.queue_item_id is None:
            return False
        item = self.db.get(DeepAnalysisQueueItem, evidence.queue_item_id)
        if item is None:
            return False
        validation = self._anchor_validation_for_evidence(evidence, item)
        return validation.anchor_validation_status != "unknown" and not validation.is_valid

    def _prefer_card_keeper(self, left: HighlightCard, right: HighlightCard) -> HighlightCard:
        if left.is_user_edited and not right.is_user_edited:
            return left
        if right.is_user_edited and not left.is_user_edited:
            return right
        return left if left.id < right.id else right

    def _narrative_meta_for_card(
        self,
        card: HighlightCard,
        evidence: Optional[StrongEvidence],
        *,
        context_preview: dict,
    ) -> dict:
        if evidence is None:
            return {
                "risk_note": "这是普通引用/背景引用素材卡片，默认不纳入报告，建议人工复核。",
                "technical_terms_used": [],
                "evidence_basis": "",
                "evidence_quote": card.evidence_quote or "",
                "evidence_context": context_preview.get("display_context") or card.evidence_quote or "",
                "key_phrases": context_preview.get("highlight_terms", []),
                "judgment_label": self._card_type_label(card.card_type),
                "why_this_judgment": "该素材卡片来自已分析论文的引用记录或弱发现，目前缺少可提升为强证据的明确正文评价，因此默认作为待复核报告素材。",
                "copy_ready_statement": card.narrative_zh or card.body_markdown or "该论文引用了目标论文，但当前系统未识别出明确强证据，建议作为待复核引用材料处理。",
                "confidence": "low",
                "target_reference_marker": "",
                "citation_text_contains_target_marker": False,
                "citation_text_contains_other_marker": False,
                "anchor_validation_status": "unknown",
                "anchor_validation_reason": "no_strong_evidence_anchor",
                "matched_template_ids": self._load_json_list(card.matched_template_ids_json),
                "matched_template_names": self._load_json_list(card.matched_template_names),
                "template_match_reason": card.template_match_reason or "",
                "template_satisfied": card.template_satisfied,
                "template_failure_reason": card.template_failure_reason or "",
            }
        item = self.db.get(DeepAnalysisQueueItem, evidence.queue_item_id) if evidence.queue_item_id else None
        narrative = generate_impact_narrative(
            evidence=evidence,
            item=item,
            card_type=card.card_type,
            context_preview=context_preview,
            notable_author=self._find_notable_author_for_item(item.id, item.citing_authors_json) if item else None,
        ) if item else {}
        direct_evidence = self._template_direct_evidence_for_strong_evidence(evidence)
        if direct_evidence:
            model_reason = str(
                direct_evidence.get("why_this_judgment_zh") or ""
            ).strip()
            model_evaluation = str(
                direct_evidence.get("copy_ready_zh") or ""
            ).strip()
            if model_reason:
                narrative["why_this_judgment"] = model_reason
                narrative["judgment_basis_zh"] = model_reason
            if model_evaluation:
                narrative["narrative_zh"] = model_evaluation
                narrative["evidence_claim_zh"] = model_evaluation
                narrative["copy_ready_statement"] = model_evaluation
                narrative["copy_ready_statement_zh"] = model_evaluation
        anchor_validation = self._anchor_validation_for_evidence(evidence, item) if item else None
        return {
            "risk_note": narrative.get("risk_note", ""),
            "technical_terms_used": narrative.get("technical_terms_used", []),
            "evidence_basis": narrative.get("evidence_basis", ""),
            "evidence_quote": narrative.get("evidence_quote", card.evidence_quote or ""),
            "evidence_context": narrative.get("evidence_context", context_preview.get("display_context") or card.evidence_quote or ""),
            "key_phrases": narrative.get("key_phrases", []),
            "judgment_label": narrative.get("judgment_label", self._card_type_label(card.card_type)),
            "why_this_judgment": narrative.get("why_this_judgment", ""),
            "copy_ready_statement": narrative.get("copy_ready_statement", card.narrative_zh or card.body_markdown or ""),
            "evidence_claim_zh": narrative.get("evidence_claim_zh", ""),
            "judgment_basis_zh": narrative.get("judgment_basis_zh", ""),
            "limitation_zh": narrative.get("limitation_zh", ""),
            "copy_ready_statement_zh": narrative.get("copy_ready_statement_zh", ""),
            "confidence_level": narrative.get("confidence_level", narrative.get("confidence", "")),
            "report_recommendation": narrative.get("report_recommendation", ""),
            "confidence": narrative.get("confidence", ""),
            "target_reference_marker": anchor_validation.target_reference_marker if anchor_validation else "",
            "citation_text_contains_target_marker": anchor_validation.citation_text_contains_target_marker if anchor_validation else False,
            "citation_text_contains_other_marker": anchor_validation.citation_text_contains_other_marker if anchor_validation else False,
            "anchor_validation_status": anchor_validation.anchor_validation_status if anchor_validation else "unknown",
            "anchor_validation_reason": anchor_validation.anchor_validation_reason if anchor_validation else "no_anchor_validation",
            "matched_template_ids": self._load_json_list(card.matched_template_ids_json)
            or self._load_json_list(evidence.matched_template_ids_json),
            "matched_template_names": self._load_json_list(card.matched_template_names),
            "template_match_reason": card.template_match_reason or evidence.template_match_reason or "",
            "template_satisfied": card.template_satisfied
            if card.template_satisfied is not None
            else evidence.template_satisfied,
            "template_failure_reason": card.template_failure_reason or evidence.template_failure_reason or "",
        }

    def _template_direct_evidence_for_strong_evidence(
        self,
        evidence: StrongEvidence,
    ) -> dict:
        result = self.db.get(FulltextAnalysisResult, evidence.fulltext_result_id)
        if result is None or result.analysis_scope != "fulltext_template_direct":
            return {}
        payload = self._load_json(result.parsed_result_json)
        evidences = payload.get("evidences")
        if not isinstance(evidences, list):
            return {}
        index = evidence.span_index
        if isinstance(index, int) and 0 <= index < len(evidences):
            direct_evidence = evidences[index]
            return direct_evidence if isinstance(direct_evidence, dict) else {}
        return {}

    def _anchor_validation_for_evidence(self, evidence: StrongEvidence, item: DeepAnalysisQueueItem):
        result = self.db.get(FulltextAnalysisResult, evidence.fulltext_result_id)
        diagnostics = self._load_json(result.candidate_spans_json if result else None)
        return validate_citation_target_anchor(
            citation_text=evidence.citation_text or "",
            target_reference_marker=diagnostics.get("target_reference_marker"),
            cited_paper_title=item.cited_paper_title,
            cited_authors_json=item.cited_authors_json,
        )

    def _add_title_slide(self, presentation: Presentation, session_id: int, cards: List[HighlightCard]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_title_bar(slide, "学术引用证据分析 - 同行评价", "报告总览")
        self._add_textbox(
            slide,
            left=0.9,
            top=1.6,
            width=11.3,
            height=3.0,
            text="学术影响力亮点评价汇报",
            font_size=28,
            bold=True,
            color=RGBColor(17, 24, 39),
        )
        overview = "\n".join(
            [
                f"学者分析会话：{session_id}",
                f"纳入报告卡片数：{len(cards)}",
                f"强证据卡片数：{sum(1 for card in cards if card.evidence_strength == 'strong')}",
                f"普通/待复核卡片数：{sum(1 for card in cards if card.evidence_strength != 'strong')}",
            ]
        )
        self._add_textbox(
            slide,
            left=0.9,
            top=2.6,
            width=5.8,
            height=2.0,
            text=overview,
            font_size=18,
            color=RGBColor(55, 65, 81),
        )

    def _add_summary_slide(self, presentation: Presentation, cards: List[HighlightCard]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_title_bar(slide, "学术引用证据分析 - 同行评价", "报告摘要")
        summary_lines = [
            f"纳入报告卡片数：{len(cards)}",
            f"理论基础 / 方法采用：{sum(1 for card in cards if card.card_type in {'theoretical_foundation', 'method_foundation'})}",
            f"代表性相关工作：{sum(1 for card in cards if card.card_type == 'representative_work')}",
            f"局限性反馈：{sum(1 for card in cards if card.card_type == 'limitation_or_negative')}",
        ]
        self._add_textbox(
            slide,
            left=0.9,
            top=1.5,
            width=5.8,
            height=4.6,
            text="\n".join(summary_lines),
            font_size=20,
            color=RGBColor(31, 41, 55),
        )
        self._add_textbox(
            slide,
            left=7.0,
            top=1.5,
            width=5.0,
            height=4.6,
            text="本页用于概览，不展示长 quote。每张卡片独立成页，quote 区块使用高亮 PNG，避免网页字段式堆叠。",
            font_size=18,
            color=RGBColor(75, 85, 99),
        )

    def _add_card_slide(self, presentation: Presentation, card: HighlightCard) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        label = self._card_type_label(card.card_type)
        self._add_title_bar(slide, "学术引用证据分析 - 同行评价", label)
        context_preview = self._context_preview_for_card(card)
        evidence = self.db.get(StrongEvidence, card.strong_evidence_id) if card.strong_evidence_id else None
        narrative_meta = self._narrative_meta_for_card(card, evidence, context_preview=context_preview)

        quote_text = self._ppt_quote_text(card, context_preview)
        quote_image = self._render_quote_image(
            session_id=card.scholar_session_id,
            card_id=card.id or 0,
            quote_text=quote_text,
            target_marker=context_preview.get("target_reference_marker") or "",
            context_preview=context_preview,
            evidence=evidence,
        )

        self._add_textbox(
            slide,
            left=0.7,
            top=1.1,
            width=6.5,
            height=0.8,
            text=self._sanitize_ppt_text(card.source_citing_paper_title, 180),
            font_size=22,
            bold=True,
            color=RGBColor(17, 24, 39),
        )
        venue_text = " / ".join(
            part for part in [card.venue or "", str(self._queue_item_year(card) or "")] if part
        )
        self._add_textbox(
            slide,
            left=0.7,
            top=1.8,
            width=4.5,
            height=0.4,
            text=venue_text,
            font_size=13,
            color=RGBColor(75, 85, 99),
        )
        slide.shapes.add_picture(str(quote_image), Inches(0.7), Inches(2.2), width=Inches(6.1))

        narrative_text = self._sanitize_ppt_text(card.narrative_zh or card.body_markdown, 520)
        self._add_textbox(
            slide,
            left=7.1,
            top=1.25,
            width=5.3,
            height=2.55,
            text=narrative_text,
            font_size=18,
            color=RGBColor(17, 24, 39),
        )
        metadata_lines = [
            f"evidence_type：{card.aspect or ''}",
            f"stance：{card.stance or ''}",
            f"score：{card.score or ''}",
        ]
        if card.notable_author_name:
            metadata_lines.insert(0, f"重要作者：{card.notable_author_name}")
        if card.fellow_status and card.fellow_status != "unknown":
            metadata_lines.insert(1 if card.notable_author_name else 0, f"荣誉身份：{card.fellow_status}")
        if narrative_meta.get("risk_note"):
            metadata_lines.append(f"风险提示：{narrative_meta['risk_note']}")
        self._add_textbox(
            slide,
            left=7.1,
            top=4.05,
            width=5.3,
            height=1.6,
            text="\n".join(metadata_lines),
            font_size=13,
            color=RGBColor(55, 65, 81),
        )
        bottom_text = self._sanitize_ppt_text(
            f"原文引句：{card.evidence_quote}",
            340,
        )
        self._add_textbox(
            slide,
            left=0.7,
            top=6.55,
            width=11.8,
            height=0.45,
            text=bottom_text,
            font_size=11,
            color=RGBColor(107, 114, 128),
        )

    def _add_title_bar(self, slide, report_title: str, card_label: str) -> None:
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(0.78),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(15, 35, 80)
        bar.line.color.rgb = RGBColor(15, 35, 80)
        self._add_textbox(
            slide,
            left=0.45,
            top=0.1,
            width=8.6,
            height=0.42,
            text=report_title,
            font_size=22,
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(10.7),
            Inches(0.12),
            Inches(2.1),
            Inches(0.42),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = self._card_type_color(card_label)
        badge.line.color.rgb = self._card_type_color(card_label)
        self._add_textbox(
            slide,
            left=10.8,
            top=0.16,
            width=1.9,
            height=0.28,
            text=card_label,
            font_size=14,
            bold=True,
            color=RGBColor(255, 255, 255),
            align=PP_ALIGN.CENTER,
        )

    def _add_textbox(
        self,
        slide,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int,
        color: RGBColor,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
    ) -> None:
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(2)
        text_frame.margin_right = Pt(2)
        text_frame.margin_top = Pt(1)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _render_quote_image(
        self,
        *,
        session_id: int,
        card_id: int,
        quote_text: str,
        target_marker: str,
        context_preview: dict,
        evidence: Optional[StrongEvidence],
    ) -> Path:
        output_path = self._session_export_dir(session_id) / f"quote_{card_id or 'card'}.png"
        highlights = [target_marker] if target_marker else []
        highlights.extend(self._load_json_list(evidence.highlight_keywords_json) if evidence else [])
        highlights.extend(str(term) for term in context_preview.get("highlight_terms", []) if str(term).strip())
        self.render_highlighted_quote_png(
            text=quote_text,
            highlights=highlights,
            output_path=output_path,
            width_px=1080,
            font_size=26,
        )
        return output_path

    def render_highlighted_quote_png(
        self,
        *,
        text: str,
        highlights: List[str],
        output_path: Path,
        width_px: int,
        font_size: int,
    ) -> None:
        font = self._quote_font(font_size)
        padding = 36
        line_spacing = 12
        lines = self._wrap_text_lines(text, font, width_px - (padding * 2))
        line_height = font_size + line_spacing
        height_px = max(260, padding * 2 + line_height * len(lines) + 24)
        image = Image.new("RGB", (width_px, height_px), (248, 250, 252))
        draw = ImageDraw.Draw(image)
        y = padding
        highlight_tokens = [token for token in highlights if token]
        for line in lines:
            lowered = line.lower()
            for token in highlight_tokens:
                token_lower = token.lower()
                start = lowered.find(token_lower)
                if start < 0:
                    continue
                prefix = line[:start]
                match_text = line[start:start + len(token)]
                x0 = padding + self._text_width(draw, prefix, font)
                x1 = x0 + self._text_width(draw, match_text, font)
                draw.rounded_rectangle((x0 - 4, y - 2, x1 + 4, y + font_size + 4), radius=6, fill=(255, 243, 153))
            draw.text((padding, y), line, fill=(17, 24, 39), font=font)
            y += line_height
        image.save(output_path)

    def _wrap_text_lines(self, text: str, font, max_width: int) -> List[str]:
        tokens = re.split(r"(\s+)", self._sanitize_ppt_text(text, 900))
        lines = []
        current = ""
        measure = ImageDraw.Draw(Image.new("RGB", (10, 10)))
        for token in tokens:
            candidate = f"{current}{token}" if current else token
            if candidate and self._text_width(measure, candidate, font) <= max_width:
                current = candidate
                continue
            if current:
                lines.append(current.strip())
                current = token.lstrip()
            else:
                lines.append(token.strip())
                current = ""
        if current.strip():
            lines.append(current.strip())
        return lines[:10]

    def _text_width(self, draw: ImageDraw.ImageDraw, text: str, font) -> int:
        if not text:
            return 0
        bbox = draw.textbbox((0, 0), text, font=font)
        return bbox[2] - bbox[0]

    def _quote_font(self, font_size: int):
        for candidate in [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
        ]:
            if Path(candidate).exists():
                try:
                    return ImageFont.truetype(candidate, font_size)
                except Exception:
                    continue
        return ImageFont.load_default()

    def _ppt_quote_text(self, card: HighlightCard, context_preview: dict) -> str:
        section = context_preview.get("section_heading") or ""
        marker = context_preview.get("target_reference_marker") or ""
        base = context_preview.get("display_context") or context_preview.get("citation_sentence") or card.evidence_quote
        if len(base) > 620:
            base = base[:620].rsplit(" ", 1)[0] + " ..."
        prefix = " | ".join(part for part in [section, marker] if part)
        return f"{prefix}\n{base}".strip()

    def _queue_item_year(self, card: HighlightCard) -> Optional[int]:
        if not card.strong_evidence_id:
            return None
        evidence = self.db.get(StrongEvidence, card.strong_evidence_id)
        if evidence is None or evidence.queue_item_id is None:
            return None
        item = self.db.get(DeepAnalysisQueueItem, evidence.queue_item_id)
        return item.year if item else None

    def _card_type_label(self, card_type: str) -> str:
        labels = {
            "theoretical_foundation": "理论基础",
            "method_foundation": "方法采用",
            "representative_work": "代表性相关工作",
            "ordinary_citation": "普通引用",
            "background_reference": "背景引用",
            "citation_only": "仅引用记录",
            "weak_mention": "弱证据",
            "limitation_or_negative": "局限性反馈",
            "detailed_comparison": "详细对比",
            "baseline_or_benchmark": "基线/Benchmark",
            "positive_evaluation": "正向评价",
            "neutral_evaluation": "中性评价",
            "first_or_seminal_claim": "首次/开创性",
        }
        return labels.get(card_type, card_type)

    def _card_type_color(self, card_label: str) -> RGBColor:
        if "理论基础" in card_label or "方法采用" in card_label:
            return RGBColor(29, 78, 216)
        if "代表性" in card_label or "普通引用" in card_label:
            return RGBColor(75, 85, 99)
        if "局限性" in card_label:
            return RGBColor(194, 65, 12)
        return RGBColor(79, 70, 229)


def get_highlight_card_service(db: Session = Depends(get_db)) -> HighlightCardService:
    return HighlightCardService(db)


def _escape_xml(value: str) -> str:
    return (
        str(value or "")
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def _content_types(slide_count: int) -> str:
    slide_overrides = "".join(
        f'<Override PartName="/ppt/slides/slide{index}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for index in range(1, slide_count + 1)
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
        '<Override PartName="/ppt/presProps.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presProps+xml"/>'
        '<Override PartName="/ppt/viewProps.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.viewProps+xml"/>'
        '<Override PartName="/ppt/tableStyles.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml"/>'
        '<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>'
        '<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
        '<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>'
        '<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>'
        '<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>'
        f"{slide_overrides}"
        "</Types>"
    )


def _root_rels() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
        '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>'
        '<Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>'
        "</Relationships>"
    )


def _presentation_xml(slide_count: int) -> str:
    slide_ids = "".join(
        f'<p:sldId id="{256 + index}" r:id="rId{index + 2}"/>'
        for index in range(slide_count)
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>'
        f"<p:sldIdLst>{slide_ids}</p:sldIdLst>"
        '<p:sldSz cx="9144000" cy="5143500"/>'
        '<p:notesSz cx="6858000" cy="9144000"/>'
        "</p:presentation>"
    )


def _presentation_rels(slide_count: int) -> str:
    slide_rels = "".join(
        f'<Relationship Id="rId{index + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{index}.xml"/>'
        for index in range(1, slide_count + 1)
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>'
        f"{slide_rels}"
        "</Relationships>"
    )


def _slide_xml(title: str, body: str) -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr/>'
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        '<p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>'
        f"{_escape_xml(title)}"
        '</a:t></a:r></a:p></p:txBody></p:sp>'
        '<p:sp><p:nvSpPr><p:cNvPr id="3" name="Content"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        '<p:spPr/><p:txBody><a:bodyPr wrap="square"/><a:lstStyle/><a:p><a:r><a:t>'
        f"{_escape_xml(body)}"
        '</a:t></a:r></a:p></p:txBody></p:sp>'
        '</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'
    )


def _core_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
        'xmlns:dc="http://purl.org/dc/elements/1.1/" '
        'xmlns:dcterms="http://purl.org/dc/terms/" '
        'xmlns:dcmitype="http://purl.org/dc/dcmitype/" '
        'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
        '<dc:title>Academic Impact Report</dc:title>'
        '<dc:creator>academic_impact_app</dc:creator>'
        "</cp:coreProperties>"
    )


def _build_slide_body(lines: List[str]) -> str:
    return "\n".join(line for line in lines if line)


_APP_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
    'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
    '<Application>academic_impact_app</Application>'
    "</Properties>"
)
_PRES_PROPS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:presentationPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
)
_VIEW_PROPS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:viewPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
)
_TABLE_STYLES_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<a:tblStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" def=""/>'
)
_THEME_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Office Theme">'
    '<a:themeElements>'
    '<a:clrScheme name="Office"><a:dk1><a:srgbClr val="000000"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1></a:clrScheme>'
    '<a:fontScheme name="Office"><a:majorFont/><a:minorFont/></a:fontScheme>'
    '<a:fmtScheme name="Office"><a:fillStyleLst/><a:lnStyleLst/><a:effectStyleLst/><a:bgFillStyleLst/></a:fmtScheme>'
    '</a:themeElements></a:theme>'
)
_SLIDE_MASTER_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    '<p:cSld name="Master"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>'
    '<p:clrMap accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" bg1="lt1" bg2="lt2" folHlink="folHlink" hlink="hlink" tx1="dk1" tx2="dk2"/>'
    '<p:sldLayoutIdLst><p:sldLayoutId id="1" r:id="rId1"/></p:sldLayoutIdLst>'
    "</p:sldMaster>"
)
_SLIDE_MASTER_RELS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
    '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>'
    "</Relationships>"
)
_SLIDE_LAYOUT_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="titleOnly" preserve="1">'
    '<p:cSld name="Layout"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>'
    '<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>'
    "</p:sldLayout>"
)
_SLIDE_LAYOUT_RELS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
)
_SLIDE_RELS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
    "</Relationships>"
)
